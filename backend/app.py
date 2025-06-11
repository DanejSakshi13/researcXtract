


























# TRYING img extraction PERFECT WORKING

# from flask import Flask, request, jsonify, send_file
# from flask_cors import CORS
# import pdfplumber
# import base64
# import os
# import json
# import google.generativeai as genai
# from dotenv import load_dotenv
# import re
# import time
# from google.api_core.exceptions import ResourceExhausted
# import requests
# from sentence_transformers import SentenceTransformer
# from numpy import dot
# from numpy.linalg import norm
# import imghdr
# import pymongo
# from bson.objectid import ObjectId
# import logging
# from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity, set_access_cookies, unset_jwt_cookies
# from reportlab.lib.pagesizes import letter
# from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
# from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
# from reportlab.lib import colors
# from io import BytesIO
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN

# app = Flask(__name__)
# CORS(app, resources={r"/api/*": {"origins": "http://localhost:5173", "methods": ["GET", "POST", "OPTIONS", "DELETE"], "allow_headers": ["Content-Type", "Authorization"]}})

# # Setup logging
# logging.basicConfig(level=logging.INFO)
# logger = logging.getLogger(__name__)

# # Load environment variables
# load_dotenv()
# GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
# MONGO_URI = os.getenv("MONGO_URI")
# JWT_SECRET_KEY = os.getenv("JWT_SECRET_KEY", "your-secret-key")
# if not GOOGLE_API_KEY:
#     raise ValueError("GOOGLE_API_KEY not found in .env file")
# if not MONGO_URI:
#     raise ValueError("MONGO_URI not found in .env file")
# genai.configure(api_key=GOOGLE_API_KEY)

# # Initialize MongoDB
# try:
#     client = pymongo.MongoClient(MONGO_URI)
#     db = client.get_database("researchXtract")
#     analyses_collection = db.analyses
#     users_collection = db.users
#     logger.info("Connected to MongoDB Atlas")
# except Exception as e:
#     logger.error(f"Failed to connect to MongoDB: {str(e)}")
#     raise

# # JWT Configuration
# app.config["JWT_SECRET_KEY"] = JWT_SECRET_KEY
# app.config["JWT_TOKEN_LOCATION"] = ["headers", "cookies"]
# jwt = JWTManager(app)

# # Handle OPTIONS requests for CORS preflight
# @app.route('/api/<path:path>', methods=['OPTIONS'])
# def handle_options(path):
#     return '', 200

# # Initialize Gemini model
# # MODEL_NAME = "gemini-1.5-pro"
# # model = genai.GenerativeModel(MODEL_NAME)
# MODEL_NAME = "gemini-2.0-flash-lite"
# model = genai.GenerativeModel(MODEL_NAME)


# # Initialize sentence transformer for semantic similarity
# sentence_model = SentenceTransformer('all-MiniLM-L6-v2')

# # Store PDF text and chat history in memory (per session)
# pdf_context = {"text": "", "session_id": None}
# chat_history = []

# def validate_input(text):
#     return len(text) > 0 and len(text) < 250000

# def query_gemini(prompt, max_tokens=512, retries=3, delay=10):
#     for attempt in range(retries):
#         try:
#             response = model.generate_content(prompt, generation_config={"max_output_tokens": max_tokens})
#             return response.text.strip()
#         except ResourceExhausted as e:
#             if attempt < retries - 1:
#                 print(f"Quota exceeded, retrying in {delay} seconds... (Attempt {attempt + 1}/{retries})")
#                 time.sleep(delay)
#                 continue
#             raise Exception(f"Gemini API quota exceeded after {retries} retries: {str(e)}")
#         except Exception as e:
#             raise Exception(f"Gemini API request failed: {str(e)}")

# def cosine_similarity(a, b):
#     return dot(a, b) / (norm(a) * norm(b))

# # Authentication Routes
# @app.route("/api/register", methods=["POST"])
# def register():
#     logger.info("Received POST request to /api/register")
#     data = request.get_json()
#     email = data.get("email")
#     password = data.get("password")
#     name = data.get("name", "Unknown")
#     if not email or not password:
#         return jsonify({"error": "Email and password are required"}), 400
#     if users_collection.find_one({"email": email}):
#         return jsonify({"error": "Email already registered"}), 400
#     user_data = {"email": email, "password": password, "name": name, "token": None}
#     logger.info(f"Registering user: {user_data}")  # Add logging
#     users_collection.insert_one(user_data)
#     return jsonify({"message": "User registered successfully."}), 201

# @app.route("/api/login", methods=["POST"])
# def login():
#     logger.info("Received POST request to /api/login")
#     data = request.get_json()
#     email = data.get("email")
#     password = data.get("password")
#     user = users_collection.find_one({"email": email, "password": password})
#     if not user:
#         logger.error(f"Login failed for {email}: Invalid credentials")
#         return jsonify({"error": "Invalid credentials"}), 401
#     access_token = create_access_token(identity=email)
#     users_collection.update_one({"email": email}, {"$set": {"token": access_token}})
#     user_data = {"email": user["email"], "name": user.get("name", "Unknown")}
#     logger.info(f"Login successful for {email}, user data: {user_data}")
#     response = jsonify({"access_token": access_token, "user": user_data})
#     set_access_cookies(response, access_token)
#     return response, 200

# @app.route("/api/user", methods=["GET"])
# @jwt_required()
# def get_user():
#     email = get_jwt_identity()
#     user = users_collection.find_one({"email": email}, {"_id": 0, "password": 0})
#     if not user:
#         logger.error(f"User not found for email: {email}")
#         return jsonify({"error": "User not found"}), 404
#     logger.info(f"User data for {email}: {user}")
#     return jsonify({"user": user}), 200

# @app.route("/api/logout", methods=["POST"])
# @jwt_required()
# def logout():
#     logger.info("Received POST request to /api/logout")
#     email = get_jwt_identity()
#     users_collection.update_one({"email": email}, {"$set": {"token": None}})
#     response = jsonify({"message": "Logged out successfully"})
#     unset_jwt_cookies(response)
#     return response, 200

# # @app.route("/api/user-history", methods=["GET"])
# # @jwt_required()
# # def user_history():
# #     email = get_jwt_identity()
# #     try:
# #         analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
# #         for analysis in analyses:
# #             analysis["_id"] = str(analysis["_id"])
# #             logger.info(f"History item {analysis['_id']}: chat_history = {analysis.get('chat_history', [])}")
# #         return jsonify({"history": analyses}), 200
# #     except Exception as e:
# #         logger.error(f"User history error: {str(e)}")
# #         return jsonify({"error": f"Failed to retrieve user history: {str(e)}"}), 500

# @app.route("/api/user-history", methods=["GET"])
# @jwt_required()
# def user_history():
#     email = get_jwt_identity()
#     try:
#         analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
#         for analysis in analyses:
#             analysis["_id"] = str(analysis["_id"])
#             # Limit recommendations to 5 for display
#             if "recommendations" in analysis and len(analysis["recommendations"]) > 5:
#                 analysis["recommendations"] = analysis["recommendations"][:5]
#                 logger.info(f"Limited recommendations for analysis {analysis['_id']} to 5")
#             logger.info(f"History item {analysis['_id']}: chat_history = {analysis.get('chat_history', [])}")
#         return jsonify({"history": analyses}), 200
#     except Exception as e:
#         logger.error(f"User history error: {str(e)}")
#         return jsonify({"error": f"Failed to retrieve user history: {str(e)}"}), 500

# @app.route("/api/analyze-pdf", methods=["POST"])
# @jwt_required()
# def analyze_pdf():
#     if "file" not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files["file"]
#     word_limit = request.form.get("word_limit", 150)
#     if not file.filename.endswith(".pdf"):
#         return jsonify({"error": "Only PDF files are supported"}), 400

#     try:
#         with pdfplumber.open(file) as pdf:
#             text = "".join(page.extract_text() or "" for page in pdf.pages)
#         print(f"Extracted text length: {len(text)}")
#         print(f"First 500 chars of extracted text: {text[:500]}")
#         if not validate_input(text):
#             return jsonify({"error": f"Invalid input: text length {len(text)} (must be 1-249999 characters)"}), 400

#         global pdf_context, chat_history
#         session_id = str(time.time())
#         pdf_context = {"text": text, "session_id": session_id}
#         chat_history = []

#         prompt_text = text[:45000]
#         if len(text) > 45000:
#             ref_start = text.rfind("REFERENCES") or text.rfind("References")
#             if ref_start != -1:
#                 ref_text = text[ref_start:ref_start+20000]
#                 prompt_text = text[:25000] + "\n" + ref_text
#                 prompt_text = prompt_text[:45000]
#             print(f"Text truncated to {len(prompt_text)} chars, References included: {ref_start != -1}")

#         prompt = f"""
#         Analyze the following research paper text and extract the specified fields.
#         Return the response as a JSON object with the following structure:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "summary": "string (summarize in {word_limit} words)",
#             "keywords": ["string", ...],
#             "citations": ["string", ...]
#         }}
#         - For "title", extract the exact title of the paper, typically found at the top of the first page.
#         - For "authors", list all author names as they appear (e.g., "Christian Szegedy").
#         - For "summary", provide a concise summary of the paper's abstract and introduction in {word_limit} words, formatted as concise bullet points.
#         - For "keywords", extract all keywords listed in the paper's "Keywords" section. If no "Keywords" section exists, infer 5–10 keywords from the abstract and introduction, focusing on technical terms and concepts central to the paper (e.g., "Inception", "convolutional neural networks").
#         - For "citations", Include the numbering (e.g., '[1]', '[2]') as it appears in the paper extract ALL citations listed in the paper's "References" or bibliography section, preserving their exact formatting as they appear (e.g., "S. Arora, A. Bhaskara, ..."). Include all citations, handling LaTeX formatting, special characters, and metadata (e.g., DOIs, URLs). Do not modify or reformat the citations.
#         - If a field cannot be extracted, use appropriate defaults (e.g., "Untitled" for title, [] for lists, "No summary available" for summary).
#         - Ensure the response is valid JSON.
#         Text (up to 45,000 characters):
#         {prompt_text}
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         print(f"Raw Gemini response: {response_text}")

#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             result = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             print(f"JSON parsing error: {str(e)}")
#             return jsonify({"error": f"Failed to parse Gemini response as JSON: {str(e)}"}), 500

#         analysis_data = {
#             "title": result.get("title", "Untitled"),
#             "authors": result.get("authors", ["No authors found"]),
#             "summary": result.get("summary", "No summary available"),
#             "keywords": result.get("keywords", []),
#             "citations": result.get("citations", []),
#             "session_id": session_id,
#             "text": text  # Include the full extracted text
#         }
#         print(f"Extracted keywords: {analysis_data['keywords']}")
#         print(f"Extracted citations: {analysis_data['citations']}")
#         return jsonify(analysis_data)
#     except Exception as e:
#         print(f"Analyze PDF error: {str(e)}")
#         return jsonify({"error": f"PDF analysis failed: {str(e)}"}), 500

# @app.route("/api/summarize-section", methods=["POST"])
# @jwt_required()
# def summarize_section():
#     data = request.get_json()
#     section = data.get("section", "")
#     session_id = data.get("session_id", "")
#     word_limit = 150

#     if not section:
#         return jsonify({"error": "No section provided"}), 400
#     if not session_id:
#         return jsonify({"error": "Missing session ID"}), 400

#     try:
#         # Retrieve the analysis document from MongoDB
#         email = get_jwt_identity()
#         analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
#         if not analysis or not analysis.get("analysis_data", {}).get("text"):
#             return jsonify({"error": "No PDF context available for this session ID."}), 400

#         prompt_text = analysis["analysis_data"]["text"][:45000]
#         print(f"Summarizing section: {section}, Context length: {len(prompt_text)} chars")

#         prompt = f"""
#         Analyze the following research paper text and provide a concise summary of the specified section in {word_limit} words.
#         - Section to summarize: "{section}"
#         - If the section is "Entire Paper", summarize the abstract and introduction (as in the default summary).
#         - For other sections (e.g., "Introduction", "Literature Review", "Methodology", "Results", "Conclusion"), identify the section by its title or content (e.g., "Related Work" for "Literature Review", "Experiments" for "Results").
#         - If the section is not found, return: "This section is not available in the paper."
#         - Return the response as a JSON object:
#         {{
#             "section": "{section}",
#             "summary": "string (summarize in {word_limit} words)"
#         }}
#         - Ensure the summary is concise, accurate, and based only on the provided text formatted as concise bullet points.
#         - Ensure the response is valid JSON.
#         Text (up to 45,000 characters):
#         {prompt_text}
#         """
#         response_text = query_gemini(prompt, max_tokens=1000)
#         print(f"Raw Gemini response for section {section}: {response_text}")

#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             result = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             print(f"JSON parsing error for section {section}: {str(e)}")
#             return jsonify({"error": f"Failed to parse Gemini response as JSON: {str(e)}"}), 500

#         response_data = {
#             "section": result.get("section", section),
#             "summary": result.get("summary", "This section is not available in the paper.")
#         }
#         print(f"Section summary for {section}: {response_data['summary'][:100]}...")

#         # Update the analysis document with the new section summary
#         analyses_collection.update_one(
#             {"session_id": session_id, "email": email},
#             {"$set": {f"analysis_data.sectionSummaries.{section}": response_data["summary"]}}
#         )

#         return jsonify(response_data)
#     except Exception as e:
#         print(f"Section summary error for {section}: {str(e)}")
#         return jsonify({"error": f"Section summary failed: {str(e)}"}), 500

# # @app.route("/api/extract-images", methods=["POST"])
# # @jwt_required()
# # def extract_images():
# #     if "file" not in request.files:
# #         return jsonify({"error": "No file provided"}), 400
# #     file = request.files["file"]
# #     if not file.filename.endswith(".pdf"):
# #         return jsonify({"error": "Only PDF files are supported"}), 400

# #     try:
# #         visuals = []
# #         with pdfplumber.open(file) as pdf:
# #             for page_num, page in enumerate(pdf.pages, 1):
# #                 try:
# #                     # Get images from page.objects for reliability
# #                     images = [obj for obj in page.objects.get('image', []) if isinstance(obj, dict) and 'stream' in obj]
# #                     logger.info(f"Page {page_num}: Found {len(images)} image objects")
# #                 except Exception as e:
# #                     logger.error(f"Page {page_num}: Failed to access page objects: {str(e)}")
# #                     continue

# #                 for img_idx, img in enumerate(images):
# #                     try:
# #                         # Validate stream
# #                         if not hasattr(img['stream'], 'get_data'):
# #                             logger.warning(f"Page {page_num}, Image {img_idx}: No valid stream data")
# #                             continue

# #                         img_data = img['stream'].get_data()
# #                         if not img_data or len(img_data) < 100:
# #                             logger.warning(f"Page {page_num}, Image {img_idx}: Image data too small or empty")
# #                             continue

# #                         # Validate image type
# #                         img_type = imghdr.what(None, img_data)
# #                         if img_type not in ['jpeg', 'png']:
# #                             logger.warning(f"Page {page_num}, Image {img_idx}: Unsupported image type {img_type or 'unknown'}")
# #                             continue

# #                         # Encode to Base64
# #                         base64_data = base64.b64encode(img_data).decode("utf-8")
# #                         if not base64_data:
# #                             logger.warning(f"Page {page_num}, Image {img_idx}: Failed to encode image to Base64")
# #                             continue

# #                         # Analyze image with Gemini to classify type
# #                         image_type = "unknown"
# #                         try:
# #                             prompt = """
# #                             Analyze the provided image and classify it as one of the following types:
# #                             - Flowchart: Diagrams with arrows and boxes showing a process or workflow.
# #                             - Figure: General scientific illustrations (e.g., graphs, charts, diagrams).
# #                             - Table: Grid-like structures with rows and columns.
# #                             - Graph: Visualizations like line graphs, bar charts, or scatter plots.
# #                             - Other: Any image that doesn't fit the above categories.
# #                             Return only the type (e.g., "Flowchart", "Figure") as a single word.
# #                             """
# #                             image_data = {"type": img_type, "data": img_data}
# #                             image_type = query_gemini(prompt, max_tokens=50, image_data=image_data)
# #                             logger.info(f"Page {page_num}, Image {img_idx}: Gemini classified as {image_type}")
# #                         except Exception as e:
# #                             logger.warning(f"Page {page_num}, Image {img_idx}: Gemini classification failed: {str(e)}")
# #                             image_type = "unknown"

# #                         visuals.append({
# #                             "name": f"image_page{page_num}_{img_idx}",
# #                             "image": base64_data,
# #                             "type": img_type,
# #                             "image_type": image_type  # New field for Gemini classification
# #                         })
# #                         logger.info(f"Page {page_num}, Image {img_idx}: Successfully extracted {img_type} image, Base64 length: {len(base64_data)}, Classified as: {image_type}")
# #                     except Exception as e:
# #                         logger.error(f"Page {page_num}, Image {img_idx}: Extraction failed: {str(e)}")
# #                         continue

# #         logger.info(f"Total valid images extracted: {len(visuals)}")
# #         return jsonify({"visuals": visuals})
# #     except Exception as e:
# #         logger.error(f"Image extraction error: {str(e)}")
# #         return jsonify({"error": f"Image extraction failed: {str(e)}"}), 500



# @app.route("/api/extract-tables", methods=["POST"])
# @jwt_required()
# def extract_tables():
#     if "file" not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files["file"]
#     if not file.filename.endswith(".pdf"):
#         return jsonify({"error": "Only PDF files are supported"}), 400

#     try:
#         tables = []
#         with pdfplumber.open(file) as pdf:
#             for page_num, page in enumerate(pdf.pages, 1):
#                 page_text = page.extract_text() or ""
#                 page_tables = page.extract_tables()
#                 print(f"Page {page_num}: Found {len(page_tables)} tables")
#                 for table_idx, table in enumerate(page_tables):
#                     try:
#                         cleaned_table = [[cell or "" for cell in row] for row in table if any(cell for cell in row)]
#                         if not cleaned_table or len(cleaned_table) < 2:
#                             print(f"Page {page_num}, Table {table_idx}: Skipped, empty or invalid")
#                             continue

#                         caption = "No caption available"
#                         caption_match = re.search(r'(?i)Table\s+\d+\s*[:\.\-\s]*(.*?)(?=\n|$)', page_text, re.MULTILINE)
#                         if caption_match:
#                             caption = caption_match.group(0).strip()
#                         else:
#                             lines = page_text.split('\n')
#                             for line in lines:
#                                 if re.match(r'(?i)Table\s+\d+', line):
#                                     caption = line.strip()
#                                     break

#                         tables.append({
#                             "page": page_num,
#                             "caption": caption,
#                             "rows": cleaned_table
#                         })
#                         print(f"Page {page_num}, Table {table_idx}: Extracted, {len(cleaned_table)} rows, Caption: {caption}")
#                     except Exception as e:
#                         print(f"Page {page_num}, Table {table_idx}: Extraction failed: {str(e)}")
#         print(f"Total tables extracted: {len(tables)}")
#         return jsonify({"tables": tables})
#     except Exception as e:
#         print(f"Table extraction error: {str(e)}")
#         return jsonify({"error": f"Table extraction failed: {str(e)}"}), 500


# @app.route("/api/recommend", methods=["POST"])
# @jwt_required()
# def recommend():
#     data = request.get_json()
#     summary = data.get("summary", "")
#     if not summary:
#         return jsonify({"error": "No summary provided"}), 400

#     try:
#         # Construct prompt for Gemini to recommend 5 recent papers
#         prompt = f"""
#         Based on the following summary of a research paper, recommend exactly 5 recent research papers (published within the last 3 years, i.e., 2022 or later) from open-access repositories (e.g., arXiv, PubMed Central, or other public repositories). 
#         The recommendations should be closely related to the topics, methods, or findings described in the summary, such as natural language processing, PDF analysis, recommendation systems, or related fields.
#         Return the response as a JSON array of 5 objects, each with the following structure:
#         {{
#             "title": "string",
#             "published": "string (YYYY-MM-DD format)",
#             "arxiv_id": "string (or other identifier if not from arXiv)",
#             "link": "string (URL to the paper)"
#         }}
#         - Ensure the papers are from open-access sources.
#         - Prioritize recent papers (2022 or later) to ensure relevance and timeliness.
#         - Provide accurate and relevant recommendations based on your knowledge.
#         - If exact publication dates or IDs are unavailable, use reasonable estimates (e.g., "2023-01-01" for recent papers) and note any assumptions.
#         - Ensure the response is valid JSON and contains exactly 5 recommendations.
#         Summary:
#         {summary[:1000]}  # Limit to 1000 characters to avoid exceeding token limits
#         """
#         response_text = query_gemini(prompt, max_tokens=1000)
#         logger.info(f"Raw Gemini response for recommendations: {response_text}")

#         # Clean up response by removing code fences
#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             recommendations = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             logger.error(f"JSON parsing error for recommendations: {str(e)}")
#             return jsonify({"error": f"Failed to parse Gemini response as JSON: {str(e)}"}), 500

#         # Validate that we have exactly 5 recommendations
#         if not isinstance(recommendations, list) or len(recommendations) != 5:
#             logger.error(f"Invalid number of recommendations: {len(recommendations)}")
#             return jsonify({"error": "Failed to generate exactly 5 recommendations"}), 500

#         # Validate the structure of each recommendation
#         for rec in recommendations:
#             if not all(key in rec for key in ["title", "published", "arxiv_id", "link"]):
#                 logger.error(f"Invalid recommendation structure: {rec}")
#                 return jsonify({"error": "Invalid recommendation structure"}), 500

#         logger.info(f"Successfully generated 5 recommendations")
#         return jsonify(recommendations)
#     except Exception as e:
#         logger.error(f"Recommendations error: {str(e)}")
#         return jsonify({"error": f"Recommendations failed: {str(e)}"}), 500

# @app.route("/api/chat", methods=["POST"])
# @jwt_required()
# def chat():
#     data = request.get_json()
#     user_message = data.get("message", "")
#     session_id = data.get("session_id", "")
#     if not user_message:
#         return jsonify({"error": "No message provided"}), 400
#     if not session_id:
#         return jsonify({"error": "Missing session ID"}), 400

#     try:
#         email = get_jwt_identity()
#         analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
#         if not analysis or not analysis.get("analysis_data", {}).get("text"):
#             return jsonify({"error": "No PDF context available for this session ID."}), 400

#         context = analysis["analysis_data"]["text"][:45000]
#         chat_history = analysis.get("chat_history", [])
#         logger.info(f"Chat history for session {session_id}: {chat_history}")
#         history_prompt = "\n".join([f"User: {msg['user']}\nAssistant: {msg['assistant']}" for msg in chat_history[-3:]])

#         prompt = f"""
#         You are an assistant that answers questions strictly based on the content of a research paper provided as context. 
#         Do not use external knowledge, make assumptions, or include information beyond the provided text.
#         Search the entire context, including all sections (e.g., Abstract, Introduction, Literature Survey, System Architecture, etc.), to find relevant information.
#         If the question cannot be answered based on the context, respond with "This information is not available in the paper."
#         Provide concise and accurate answers, directly addressing the user's question.
        
#         Context (up to 45,000 characters):
#         {context}
        
#         Recent conversation:
#         {history_prompt}
        
#         User question: {user_message}
        
#         Answer:
#         """
#         print(f"Chat prompt length: {len(prompt)} chars, Context length: {len(context)} chars")
#         response_text = query_gemini(prompt, max_tokens=500)

#         chat_history.append({"user": user_message, "assistant": response_text})
#         analyses_collection.update_one(
#             {"session_id": session_id, "email": email},
#             {"$set": {"chat_history": chat_history}},
#             upsert=True
#         )
#         logger.info(f"Updated chat history for session {session_id}: {chat_history}")

#         return jsonify({"response": response_text, "session_id": session_id})
#     except Exception as e:
#         logger.error(f"Chat error: {str(e)}")
#         return jsonify({"error": f"Chat failed: {str(e)}"}), 500


# # @app.route("/api/paper-analysis", methods=["POST"])
# # @jwt_required()
# # def save_analysis():
# #     data = request.get_json()
# #     email = get_jwt_identity()
# #     analysis_data = data.get("analysis_data")
# #     if not analysis_data:
# #         return jsonify({"error": "Missing analysis data"}), 400

# #     try:
# #         analysis_doc = {
# #             "email": email,
# #             "session_id": analysis_data.get("session_id", str(time.time())),
# #             "analysis_data": analysis_data,
# #             "tables": analysis_data.get("tables", []),
# #             "figures": analysis_data.get("figures", []),
# #             "recommendations": analysis_data.get("recommendations", []),
# #             "chat_history": analysis_data.get("messages", []),
# #             "created_at": time.time()
# #         }
# #         logger.info(f"Saving analysis for {email}, chat_history: {analysis_doc['chat_history']}")
# #         result = analyses_collection.insert_one(analysis_doc)
# #         logger.info(f"Stored analysis for {email}: {result.inserted_id}")
# #         return jsonify({"message": "Analysis saved", "id": str(result.inserted_id)})
# #     except Exception as e:
# #         logger.error(f"Save analysis error: {str(e)}")
# #         return jsonify({"error": f"Failed to save analysis: {str(e)}"}), 500

# @app.route("/api/paper-analysis", methods=["POST"])
# @jwt_required()
# def save_analysis():
#     data = request.get_json()
#     email = get_jwt_identity()
#     analysis_data = data.get("analysis_data")
#     if not analysis_data:
#         return jsonify({"error": "Missing analysis data"}), 400

#     try:
#         # Limit recommendations to 5 unique entries
#         recommendations = analysis_data.get("recommendations", [])
#         logger.info(f"Initial recommendations count: {len(recommendations)}")
        
#         # Deduplicate based on title and limit to 5
#         unique_recommendations = []
#         seen_titles = set()
#         for rec in recommendations:
#             title = rec.get("title", "")
#             if title and title not in seen_titles:
#                 seen_titles.add(title)
#                 unique_recommendations.append(rec)
#             if len(unique_recommendations) >= 5:
#                 break

#         logger.info(f"After deduplication, recommendations count: {len(unique_recommendations)}")

#         analysis_doc = {
#             "email": email,
#             "session_id": analysis_data.get("session_id", str(time.time())),
#             "analysis_data": analysis_data,
#             "tables": analysis_data.get("tables", []),
#             "figures": analysis_data.get("figures", []),
#             "recommendations": unique_recommendations,  # Use deduplicated recommendations
#             "chat_history": analysis_data.get("messages", []),
#             "created_at": time.time()
#         }
#         logger.info(f"Saving analysis for {email}, recommendations: {len(unique_recommendations)}, chat_history: {analysis_doc['chat_history']}")
#         result = analyses_collection.insert_one(analysis_doc)
#         logger.info(f"Stored analysis for {email}: {result.inserted_id}")
#         return jsonify({"message": "Analysis saved", "id": str(result.inserted_id)})
#     except Exception as e:
#         logger.error(f"Save analysis error: {str(e)}")
#         return jsonify({"error": f"Failed to save analysis: {str(e)}"}), 500

# @app.route("/api/paper-analysis/<analysis_id>", methods=["DELETE"])
# @jwt_required()
# def delete_analysis(analysis_id):
#     email = get_jwt_identity()
#     try:
#         result = analyses_collection.delete_one({"_id": ObjectId(analysis_id), "email": email})
#         if result.deleted_count == 1:
#             logger.info(f"Deleted analysis {analysis_id} for {email}")
#             return jsonify({"message": "Analysis deleted successfully"}), 200
#         else:
#             return jsonify({"error": "Analysis not found or unauthorized"}), 404
#     except Exception as e:
#         logger.error(f"Delete analysis error: {str(e)}")
#         return jsonify({"error": f"Failed to delete analysis: {str(e)}"}), 500


# @app.route("/api/user-data", methods=["GET"])
# @jwt_required()
# def user_data():
#     email = get_jwt_identity()
#     try:
#         analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
#         for analysis in analyses:
#             analysis["_id"] = str(analysis["_id"])
#         logger.info(f"Retrieved {len(analyses)} analyses for {email}")
#         return jsonify({"recent_analysis": analyses})
#     except Exception as e:
#         logger.error(f"User data error: {str(e)}")
#         return jsonify({"error": f"Failed to retrieve user data: {str(e)}"}), 500


# @app.route("/api/generate-pdf", methods=["POST"])
# @jwt_required()
# def generate_pdf():
#     try:
#         data = request.get_json()
#         analysis_data = data.get("analysis_data")
#         if not analysis_data:
#             return jsonify({"error": "Missing analysis data"}), 400

#         # Log analysis_data to debug recommendations and methodology
#         logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
#         logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

#         # Prepare content using Gemini
#         prompt = f"""
#         Generate a professional PDF report for a research paper based on the following analysis data.
#         Format the content as a structured JSON object with sections suitable for a conference presentation.
#         Include the following sections:
#         - Title: The paper's title
#         - Authors: List of authors
#         - Abstract: A concise summary (use the provided summary)
#         - Methodology: A concise summary of the paper's methodology (use sectionSummaries.Methodology if available, otherwise infer from the text)
#         - Keywords: List of keywords
#         - Tables: Full table data including captions and rows/columns (limit to 3 tables for brevity)
#         - Citations: Key citations (limit to 5 for brevity)
#         - Recommendations: All recommended papers (up to 5, include title, published date, and link)
#         For tables, include each table's caption and its full data as a list of lists (rows and columns).
#         For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
#         For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
#         Return the response as a valid JSON object:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "abstract": "string",
#             "methodology": "string",
#             "keywords": ["string", ...],
#             "tables": [
#                 {{
#                     "caption": "string",
#                     "rows": [["string", ...], ...]
#                 }},
#                 ...
#             ],
#             "citations": ["string", ...],
#             "recommendations": ["string", ...]
#         }}
#         Analysis Data:
#         {json.dumps(analysis_data, indent=2)[:10000]}  # Limit to avoid token issues
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
#         report_content = json.loads(response_text)

#         # Log report_content to debug recommendations and methodology
#         logger.info(f"report_content.recommendations: {report_content.get('recommendations', [])}")
#         logger.info(f"report_content.methodology: {report_content.get('methodology', 'Not provided')}")

#         # Create PDF using reportlab
#         buffer = BytesIO()
#         doc = SimpleDocTemplate(buffer, pagesize=letter)
#         styles = getSampleStyleSheet()
#         title_style = ParagraphStyle(
#             name='TitleStyle',
#             fontSize=16,
#             leading=20,
#             alignment=1,  # Center
#             spaceAfter=12,
#             fontName='Helvetica-Bold'
#         )
#         heading_style = ParagraphStyle(
#             name='HeadingStyle',
#             fontSize=12,
#             leading=14,
#             spaceAfter=8,
#             fontName='Helvetica-Bold'
#         )
#         normal_style = ParagraphStyle(
#             name='NormalStyle',
#             fontSize=10,
#             leading=12,
#             spaceAfter=6
#         )

#         story = []
#         # Title
#         story.append(Paragraph(report_content.get("title", "Untitled"), title_style))
#         story.append(Spacer(1, 12))

#         # Authors
#         story.append(Paragraph("Authors", heading_style))
#         authors = ", ".join(report_content.get("authors", ["No authors found"]))
#         story.append(Paragraph(authors, normal_style))
#         story.append(Spacer(1, 12))

#         # Abstract
#         story.append(Paragraph("Abstract", heading_style))
#         story.append(Paragraph(report_content.get("abstract", "No summary available"), normal_style))
#         story.append(Spacer(1, 12))

#         # Methodology
#         story.append(Paragraph("Methodology", heading_style))
#         methodology = report_content.get("methodology", "No methodology available")
#         story.append(Paragraph(methodology, normal_style))
#         story.append(Spacer(1, 12))

#         # Keywords
#         story.append(Paragraph("Keywords", heading_style))
#         keywords = ", ".join(report_content.get("keywords", []))
#         story.append(Paragraph(keywords, normal_style))
#         story.append(Spacer(1, 12))

#         # Tables
#         story.append(Paragraph("Tables", heading_style))
#         tables = report_content.get("tables", [])[:3]  # Limit to 3 tables
#         if tables:
#             for i, table in enumerate(tables, 1):
#                 caption = table.get("caption", f"Table {i}")
#                 rows = table.get("rows", [])
#                 if rows:
#                     # Convert rows to a format suitable for reportlab Table
#                     table_data = [[Paragraph(str(cell), normal_style) for cell in row] for row in rows]
#                     # Create table
#                     pdf_table = Table(table_data)
#                     pdf_table.setStyle(TableStyle([
#                         ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
#                         ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
#                         ('VALIGN', (0, 0), (-1, -1), 'TOP'),
#                         ('FONTSIZE', (0, 0), (-1, -1), 8),
#                     ]))
#                     story.append(Paragraph(caption, normal_style))
#                     story.append(pdf_table)
#                     story.append(Spacer(1, 12))
#                 else:
#                     story.append(Paragraph(f"{caption}: No data available", normal_style))
#         else:
#             story.append(Paragraph("No tables available", normal_style))
#         story.append(Spacer(1, 12))

#         # Citations
#         story.append(Paragraph("Citations", heading_style))
#         citations = report_content.get("citations", [])[:5]
#         for i, citation in enumerate(citations, 1):
#             story.append(Paragraph(f"{i}. {citation}", normal_style))
#         if not citations:
#             story.append(Paragraph("No citations available", normal_style))
#         story.append(Spacer(1, 12))

#         # Recommendations
#         story.append(Paragraph("Recommendations", heading_style))
#         recommendations = report_content.get("recommendations", [])  # No limit, use all
#         if recommendations:
#             for i, rec in enumerate(recommendations, 1):
#                 story.append(Paragraph(f"{i}. {rec}", normal_style))
#         else:
#             story.append(Paragraph("No recommendations available", normal_style))
#         story.append(Spacer(1, 12))

#         doc.build(story)
#         buffer.seek(0)
#         return send_file(
#             buffer,
#             as_attachment=True,
#             download_name=f"{report_content.get('title', 'report')}.pdf",
#             mimetype='application/pdf'
#         )
#     except Exception as e:
#         logger.error(f"PDF generation error: {str(e)}")
#         return jsonify({"error": f"Failed to generate PDF: {str(e)}"}), 500

# @app.route("/api/generate-ppt", methods=["POST"])
# @jwt_required()
# def generate_ppt():
#     try:
#         data = request.get_json()
#         analysis_data = data.get("analysis_data")
#         if not analysis_data:
#             return jsonify({"error": "Missing analysis data"}), 400

#         # Log analysis_data to debug recommendations and methodology
#         logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
#         logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

#         # Prepare content using Gemini
#         prompt = f"""
#         Generate content for a PowerPoint presentation for a research paper based on the following analysis data.
#         Format the content as a structured JSON object with sections suitable for a conference presentation.
#         Include the following sections for slides:
#         - Title: The paper's title and authors
#         - Summary: A concise summary (use the provided summary)
#         - Methodology: A concise summary of the paper's methodology (use sectionSummaries.Methodology if available, otherwise infer from the text)
#         - Keywords: List of keywords
#         - Tables: Full table data including captions and rows/columns (limit to 2 tables for brevity)
#         - Citations: Key citations (limit to 3)
#         - Recommendations: All recommended papers (up to 5, include title, published date, and link)
#         For tables, include each table's caption and its full data as a list of lists (rows and columns).
#         For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
#         For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
#         Return the response as a valid JSON object:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "summary": "string",
#             "methodology": "string",
#             "keywords": ["string", ...],
#             "tables": [
#                 {{
#                     "caption": "string",
#                     "rows": [["string", ...], ...]
#                 }},
#                 ...
#             ],
#             "citations": ["string", ...],
#             "recommendations": ["string", ...]
#         }}
#         Analysis Data:
#         {json.dumps(analysis_data, indent=2)[:10000]}  # Limit to avoid token issues
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
#         ppt_content = json.loads(response_text)

#         # Log ppt_content to debug recommendations and methodology
#         logger.info(f"ppt_content.recommendations: {ppt_content.get('recommendations', [])}")
#         logger.info(f"ppt_content.methodology: {ppt_content.get('methodology', 'Not provided')}")

#         # Create PPT using python-pptx
#         prs = Presentation()
#         slide_width = prs.slide_width
#         slide_height = prs.slide_height

#         # Title Slide
#         slide_layout = prs.slide_layouts[0]  # Title Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         subtitle = slide.placeholders[1]
#         title.text = ppt_content.get("title", "Untitled")
#         title.text_frame.paragraphs[0].font.size = Pt(32)
#         subtitle.text = ", ".join(ppt_content.get("authors", ["No authors found"]))
#         subtitle.text_frame.paragraphs[0].font.size = Pt(20)

#         # Summary Slide
#         slide_layout = prs.slide_layouts[1]  # Title and Content
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Summary"
#         content.text = ppt_content.get("summary", "No summary available")
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         # Methodology Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Methodology"
#         methodology = ppt_content.get("methodology", "No methodology available")
#         if not methodology or methodology.strip() == "":
#             methodology = "No methodology available"
#             logger.warning("Methodology field empty or missing in ppt_content")
#         content.text = methodology
#         content.text_frame.paragraphs[0].font.size = Pt(16)  # Slightly smaller font to fit 100-150 words
#         for p in content.text_frame.paragraphs:
#             p.font.size = Pt(16)  # Ensure all paragraphs are consistent

#         # Keywords Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Keywords"
#         keywords = ppt_content.get("keywords", [])
#         content.text = ", ".join(keywords) if keywords else "No keywords available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         # Tables Slide(s)
#         tables = ppt_content.get("tables", [])[:2]  # Limit to 2 tables
#         if tables:
#             for i, table in enumerate(tables, 1):
#                 slide = prs.slides.add_slide(slide_layout)
#                 title = slide.shapes.title
#                 title.text = f"Table {i}: {table.get('caption', 'Table')}"

#                 rows = table.get("rows", [])
#                 if rows:
#                     # Determine table dimensions (limit to 10 rows and 6 columns to fit slide)
#                     num_rows = min(len(rows), 10)
#                     num_cols = min(len(rows[0]), 6) if rows else 0
#                     if num_rows > 0 and num_cols > 0:
#                         # Add table to slide
#                         table_shape = slide.shapes.add_table(
#                             rows=num_rows,
#                             cols=num_cols,
#                             left=Inches(1),
#                             top=Inches(1.5),
#                             width=Inches(8),
#                             height=Inches(4)
#                         )
#                         ppt_table = table_shape.table
#                         # Populate table
#                         for row_idx in range(num_rows):
#                             for col_idx in range(num_cols):
#                                 cell = ppt_table.cell(row_idx, col_idx)
#                                 cell.text = str(rows[row_idx][col_idx])
#                                 cell.text_frame.paragraphs[0].font.size = Pt(12)
#                         # Style table
#                         ppt_table.first_row = True  # Header row
#                     else:
#                         content = slide.placeholders[1]
#                         content.text = "Table data too large or invalid"
#                         content.text_frame.paragraphs[0].font.size = Pt(18)
#                 else:
#                     content = slide.placeholders[1]
#                     content.text = "No table data available"
#                     content.text_frame.paragraphs[0].font.size = Pt(18)
#         else:
#             slide = prs.slides.add_slide(slide_layout)
#             title = slide.shapes.title
#             content = slide.placeholders[1]
#             title.text = "Tables"
#             content.text = "No tables available"
#             content.text_frame.paragraphs[0].font.size = Pt(18)

#         # Citations Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Citations"
#         citations = ppt_content.get("citations", [])[:3]
#         if citations:
#             for i, citation in enumerate(citations, 1):
#                 p = content.text_frame.add_paragraph()
#                 p.text = f"{i}. {citation}"
#                 p.font.size = Pt(16)
#         else:
#             content.text = "No citations available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         # Recommendations Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Recommendations"
#         recommendations = ppt_content.get("recommendations", [])  # No limit, use all
#         if recommendations:
#             for i, rec in enumerate(recommendations, 1):
#                 p = content.text_frame.add_paragraph()
#                 p.text = f"{i}. {rec}"
#                 p.font.size = Pt(14)  # Smaller font to fit 5 items
#                 p.level = 1  # Indent as bullet points
#         else:
#             content.text = "No recommendations available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         # Thank You Slide
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Thank You"
#         content.text = "Questions?"
#         content.text_frame.paragraphs[0].font.size = Pt(24)
#         content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

#         # Save PPT to buffer
#         buffer = BytesIO()
#         prs.save(buffer)
#         buffer.seek(0)
#         return send_file(
#             buffer,
#             as_attachment=True,
#             download_name=f"{ppt_content.get('title', 'presentation')}.pptx",
#             mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
#         )
#     except Exception as e:
#         logger.error(f"PPT generation error: {str(e)}")
#         return jsonify({"error": f"Failed to generate PPT: {str(e)}"}), 500



# @app.route("/api/convert-citations", methods=["POST"])
# @jwt_required()
# def convert_citations():
#     try:
#         data = request.get_json()
#         citations = data.get("citations", [])
#         style = data.get("style", "APA").upper()  # Default to APA, normalize to uppercase

#         if not citations:
#             return jsonify({"error": "No citations provided"}), 400
#         if style not in ["APA", "MLA", "CHICAGO", "IEEE"]:
#             return jsonify({"error": f"Invalid citation style: {style}. Supported styles: APA, MLA, Chicago, IEEE"}), 400

#         # Prepare prompt for Gemini to convert citations
#         prompt = f"""
#         Convert the following citations to the {style} citation style.
#         The input citations are raw strings extracted from a research paper's References section.
#         Parse each citation to identify key components (e.g., authors, title, journal, year, volume, pages, DOI, URL) and reformat them according to the {style} style guidelines.
#         - For APA: Follow the 7th edition (e.g., Author, A. A. (Year). Title. Journal, Volume(Issue), pages. DOI).
#         - For MLA: Follow the 9th edition (e.g., Author. "Title." Journal, Volume, Issue, Year, pages, DOI or URL).
#         - For Chicago (Author-Date): Follow the 17th edition (e.g., Author, First. Year. "Title." Journal Volume (Issue): pages. DOI).
#         - For IEEE: Follow the IEEE style (e.g., [1] A. Author, "Title," Journal, vol. Volume, no. Issue, pp. pages, Year).
#         - If a citation lacks sufficient information (e.g., missing year or title), include as much as possible and note "[Incomplete]" at the end.
#         - Preserve all available metadata (e.g., DOIs, URLs) and include them as required by the style.
#         - Return the response as a JSON array of strings, where each string is a citation formatted in the {style} style.
#         - Ensure the response is valid JSON and matches the number of input citations.
#         - Wrap the JSON in triple backticks with 'json' identifier:
#         ```json
#         ["formatted citation 1", "formatted citation 2", ...]
#         ```
#         Input Citations:
#         {json.dumps(citations, indent=2)[:10000]}  # Limit to avoid token issues
#         """
#         max_tokens = 3000  # Increased to handle larger responses
#         response_text = query_gemini(prompt, max_tokens=max_tokens)
#         logger.info(f"Raw Gemini response for {style} conversion: {response_text}")

#         # Clean the response: Extract JSON between ```json ... ``` or fallback to raw text
#         json_match = re.search(r'```json\n([\s\S]*?)\n```', response_text, re.MULTILINE)
#         if json_match:
#             response_text = json_match.group(1).strip()
#         else:
#             # Fallback: Remove any non-JSON content
#             response_text = re.sub(r'^.*?\[|\].*?$', '', response_text, flags=re.DOTALL).strip()
#             response_text = f'[{response_text}]'

#         # Attempt to parse JSON
#         try:
#             converted_citations = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             logger.error(f"JSON parsing error for citation conversion: {str(e)}")
#             logger.error(f"Problematic response_text: {response_text}")
#             # Fallback: Return original citations with error note
#             converted_citations = [f"{citation} [Conversion failed]" for citation in citations]
#             logger.warning(f"Falling back to original citations with error note")

#         # Validate response
#         if not isinstance(converted_citations, list) or len(converted_citations) != len(citations):
#             logger.error(f"Invalid converted citations: Expected {len(citations)}, got {len(converted_citations)}")
#             converted_citations = [f"{citation} [Invalid conversion]" for citation in citations]

#         logger.info(f"Successfully converted {len(converted_citations)} citations to {style} style")
#         return jsonify({"converted_citations": converted_citations})
#     except Exception as e:
#         logger.error(f"Citation conversion error: {str(e)}")
#         return jsonify({"error": f"Failed to convert citations: {str(e)}"}), 500



# if __name__ == "__main__":
#     app.run(debug=True, host="0.0.0.0", port=5000)


























































# perfect - after img extraction
# from flask import Flask, request, jsonify, send_file
# from flask_cors import CORS
# import pdfplumber
# import base64
# import os
# import json
# import google.generativeai as genai
# from dotenv import load_dotenv
# import re
# import time
# from google.api_core.exceptions import ResourceExhausted
# import requests
# from sentence_transformers import SentenceTransformer
# from numpy import dot
# from numpy.linalg import norm
# import imghdr
# import pymongo
# import scipy.stats as stats
# from bson.objectid import ObjectId
# import logging
# from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity, set_access_cookies, unset_jwt_cookies
# from reportlab.lib.pagesizes import letter
# from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
# from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
# from reportlab.lib import colors
# from io import BytesIO
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pdf2image import convert_from_bytes
# import numpy as np
# import cv2
# from PIL import Image
# import io

# app = Flask(__name__)
# CORS(app, resources={r"/api/*": {"origins": "http://localhost:5173", "methods": ["GET", "POST", "OPTIONS", "DELETE"], "allow_headers": ["Content-Type", "Authorization"]}})

# # Setup logging
# logging.basicConfig(level=logging.INFO)
# logger = logging.getLogger(__name__)

# # Load environment variables
# load_dotenv()
# GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
# MONGO_URI = os.getenv("MONGO_URI")
# JWT_SECRET_KEY = os.getenv("JWT_SECRET_KEY", "your-secret-key")
# if not GOOGLE_API_KEY:
#     raise ValueError("GOOGLE_API_KEY not found in .env file")
# if not MONGO_URI:
#     raise ValueError("MONGO_URI not found in .env file")
# genai.configure(api_key=GOOGLE_API_KEY)

# # Initialize MongoDB
# try:
#     client = pymongo.MongoClient(MONGO_URI)
#     db = client.get_database("researchXtract")
#     analyses_collection = db.analyses
#     users_collection = db.users
#     logger.info("Connected to MongoDB Atlas")
# except Exception as e:
#     logger.error(f"Failed to connect to MongoDB: {str(e)}")
#     raise

# # JWT Configuration
# app.config["JWT_SECRET_KEY"] = JWT_SECRET_KEY
# app.config["JWT_TOKEN_LOCATION"] = ["headers", "cookies"]
# jwt = JWTManager(app)

# # Handle OPTIONS requests for CORS preflight
# @app.route('/api/<path:path>', methods=['OPTIONS'])
# def handle_options(path):
#     return '', 200

# # Initialize Gemini model
# MODEL_NAME = "gemini-2.0-flash-lite"
# model = genai.GenerativeModel(MODEL_NAME)

# # Initialize sentence transformer for semantic similarity
# sentence_model = SentenceTransformer('all-MiniLM-L6-v2')

# # Store PDF text and chat history in memory (per session)
# pdf_context = {"text": "", "session_id": None}
# chat_history = []

# def validate_input(text):
#     return len(text) > 0 and len(text) < 250000

# def query_gemini(prompt, max_tokens=512, retries=3, delay=10, image_data=None):
#     for attempt in range(retries):
#         try:
#             if image_data:
#                 # Simulate image handling (Gemini API doesn't support direct image input in this context)
#                 # In a real implementation, you'd need to use a different approach
#                 logger.warning("Image data provided but not supported in this implementation.")
#             response = model.generate_content(prompt, generation_config={"max_output_tokens": max_tokens})
#             return response.text.strip()
#         except ResourceExhausted as e:
#             if attempt < retries - 1:
#                 print(f"Quota exceeded, retrying in {delay} seconds... (Attempt {attempt + 1}/{retries})")
#                 time.sleep(delay)
#                 continue
#             raise Exception(f"Quota exceeded after {retries} retries: {str(e)}")
#         except Exception as e:
#             raise Exception(f"API request failed: {str(e)}")

# def cosine_similarity(a, b):
#     return dot(a, b) / (norm(a) * norm(b))

# # Authentication Routes
# @app.route("/api/register", methods=["POST"])
# def register():
#     logger.info("Received POST request to /api/register")
#     data = request.get_json()
#     email = data.get("email")
#     password = data.get("password")
#     name = data.get("name", "Unknown")
#     if not email or not password:
#         return jsonify({"error": "Email and password are required"}), 400
#     if users_collection.find_one({"email": email}):
#         return jsonify({"error": "Email already registered"}), 400
#     user_data = {"email": email, "password": password, "name": name, "token": None}
#     logger.info(f"Registering user: {user_data}")
#     users_collection.insert_one(user_data)
#     return jsonify({"message": "User registered successfully."}), 201

# @app.route("/api/login", methods=["POST"])
# def login():
#     logger.info("Received POST request to /api/login")
#     data = request.get_json()
#     email = data.get("email")
#     password = data.get("password")
#     user = users_collection.find_one({"email": email, "password": password})
#     if not user:
#         logger.error(f"Login failed for {email}: Invalid credentials")
#         return jsonify({"error": "Invalid credentials"}), 401
#     access_token = create_access_token(identity=email)
#     users_collection.update_one({"email": email}, {"$set": {"token": access_token}})
#     user_data = {"email": user["email"], "name": user.get("name", "Unknown")}
#     logger.info(f"Login successful for {email}, user data: {user_data}")
#     response = jsonify({"access_token": access_token, "user": user_data})
#     set_access_cookies(response, access_token)
#     return response, 200

# @app.route("/api/user", methods=["GET"])
# @jwt_required()
# def get_user():
#     email = get_jwt_identity()
#     user = users_collection.find_one({"email": email}, {"_id": 0, "password": 0})
#     if not user:
#         logger.error(f"User not found for email: {email}")
#         return jsonify({"error": "User not found"}), 404
#     logger.info(f"User data for {email}: {user}")
#     return jsonify({"user": user}), 200

# @app.route("/api/logout", methods=["POST"])
# @jwt_required()
# def logout():
#     logger.info("Received POST request to /api/logout")
#     email = get_jwt_identity()
#     users_collection.update_one({"email": email}, {"$set": {"token": None}})
#     response = jsonify({"message": "Logged out successfully"})
#     unset_jwt_cookies(response)
#     return response, 200

# @app.route("/api/user-history", methods=["GET"])
# @jwt_required()
# def user_history():
#     email = get_jwt_identity()
#     try:
#         analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
#         for analysis in analyses:
#             analysis["_id"] = str(analysis["_id"])
#             if "recommendations" in analysis and len(analysis["recommendations"]) > 5:
#                 analysis["recommendations"] = analysis["recommendations"][:5]
#                 logger.info(f"Limited recommendations for analysis {analysis['_id']} to 5")
#             logger.info(f"History item {analysis['_id']}: chat_history = {analysis.get('chat_history', [])}")
#         return jsonify({"history": analyses}), 200
#     except Exception as e:
#         logger.error(f"User history error: {str(e)}")
#         return jsonify({"error": f"Failed to retrieve user history: {str(e)}"}), 500

# @app.route("/api/analyze-pdf", methods=["POST"])
# @jwt_required()
# def analyze_pdf():

#     # logger.info("Starting 150-second delay for /api/analyze-pdf")
#     # time.sleep(100)
#     # logger.info("Delay completed, proceeding with PDF analysis")

#     if "file" not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files["file"]
#     word_limit = request.form.get("word_limit", 150)
#     if not file.filename.endswith(".pdf"):
#         return jsonify({"error": "Only PDF files are supported"}), 400

#     try:
#         with pdfplumber.open(file) as pdf:
#             text = "".join(page.extract_text() or "" for page in pdf.pages)
#         print(f"Extracted text length: {len(text)}")
#         print(f"First 500 chars of extracted text: {text[:500]}")
#         if not validate_input(text):
#             return jsonify({"error": f"Invalid input: text length {len(text)} (must be 1-249999 characters)"}), 400

#         global pdf_context, chat_history
#         session_id = str(time.time())
#         pdf_context = {"text": text, "session_id": session_id}
#         chat_history = []

#         prompt_text = text[:45000]
#         if len(text) > 45000:
#             ref_start = text.rfind("REFERENCES") or text.rfind("References")
#             if ref_start != -1:
#                 ref_text = text[ref_start:ref_start+20000]
#                 prompt_text = text[:25000] + "\n" + ref_text
#                 prompt_text = prompt_text[:45000]
#             print(f"Text truncated to {len(prompt_text)} chars, References included: {ref_start != -1}")

#         prompt = f"""
#         Analyze the following research paper text and extract the specified fields.
#         Return the response as a JSON object with the following structure:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "summary": "string (summarize in {word_limit} words)",
#             "keywords": ["string", ...],
#             "citations": ["string", ...]
#         }}
#         - For "title", extract the exact title of the paper, typically found at the top of the first page.
#         - For "authors", list all author names as they appear (e.g., "Christian Szegedy").
#         - For "summary", provide a concise summary of the paper's abstract and introduction in {word_limit} words, formatted as concise bullet points.
#         - For "keywords", extract all keywords listed in the paper's "Keywords" section. If no "Keywords" section exists, infer 5–10 keywords from the abstract and introduction, focusing on technical terms and concepts central to the paper (e.g., "Inception", "convolutional neural networks").
#         - For "citations", extract ALL citations listed in the paper's "References" or bibliography section, preserving their exact formatting as they appear in the original text (e.g., "[1] S. Arora, A. Bhaskara, ..."). Include the numbering (e.g., "[1]", "[2]") as it appears in the paper, along with all other details such as author names, titles, journal names, years, DOIs, URLs, and any LaTeX formatting or special characters. Do not modify, reformat, or remove any part of the citation, including the numbering.
#         - If a field cannot be extracted, use appropriate defaults (e.g., "Untitled" for title, [] for lists, "No summary available" for summary).
#         - Ensure the response is valid JSON.
#         Text (up to 45,000 characters):
#         {prompt_text}
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         print(f"Raw response: {response_text}")

#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             result = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             print(f"JSON parsing error: {str(e)}")
#             return jsonify({"error": f"Failed to parse response as JSON: {str(e)}"}), 500

#         analysis_data = {
#             "title": result.get("title", "Untitled"),
#             "authors": result.get("authors", ["No authors found"]),
#             "summary": result.get("summary", "No summary available"),
#             "keywords": result.get("keywords", []),
#             "citations": result.get("citations", []),
#             "session_id": session_id,
#             "text": text
#         }
#         print(f"Extracted keywords: {analysis_data['keywords']}")
#         print(f"Extracted citations (with numbering check): {analysis_data['citations'][:5]}")
#         return jsonify(analysis_data)
#     except Exception as e:
#         print(f"Analyze PDF error: {str(e)}")
#         return jsonify({"error": f"PDF analysis failed: {str(e)}"}), 500

# @app.route("/api/summarize-section", methods=["POST"])
# @jwt_required()
# def summarize_section():
#     data = request.get_json()
#     section = data.get("section", "")
#     session_id = data.get("session_id", "")
#     word_limit = 150

#     if not section:
#         return jsonify({"error": "No section provided"}), 400
#     if not session_id:
#         return jsonify({"error": "Missing session ID"}), 400

#     try:
#         email = get_jwt_identity()
#         analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
#         if not analysis or not analysis.get("analysis_data", {}).get("text"):
#             return jsonify({"error": "No PDF context available for this session ID."}), 400

#         prompt_text = analysis["analysis_data"]["text"][:45000]
#         print(f"Summarizing section: {section}, Context length: {len(prompt_text)} chars")

#         prompt = f"""
#         Analyze the following research paper text and provide a concise summary of the specified section in {word_limit} words.
#         - Section to summarize: "{section}"
#         - If the section is "Entire Paper", summarize the abstract and introduction (as in the default summary).
#         - For other sections (e.g., "Introduction", "Literature Review", "Methodology", "Results", "Conclusion"), identify the section by its title or content (e.g., "Related Work" for "Literature Review", "Experiments" for "Results").
#         - If the section is not found, return: "This section is not available in the paper."
#         - Return the response as a JSON object:
#         {{
#             "section": "{section}",
#             "summary": "string (summarize in {word_limit} words)"
#         }}
#         - Ensure the summary is concise, accurate, and based only on the provided text formatted as concise bullet points.
#         - Ensure the response is valid JSON.
#         Text (up to 45,000 characters):
#         {prompt_text}
#         """
#         response_text = query_gemini(prompt, max_tokens=1000)
#         print(f"Raw response for section {section}: {response_text}")

#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             result = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             print(f"JSON parsing error for section {section}: {str(e)}")
#             return jsonify({"error": f"Failed to parse response as JSON: {str(e)}"}), 500

#         response_data = {
#             "section": result.get("section", section),
#             "summary": result.get("summary", "This section is not available in the paper.")
#         }
#         print(f"Section summary for {section}: {response_data['summary'][:100]}...")

#         analyses_collection.update_one(
#             {"session_id": session_id, "email": email},
#             {"$set": {f"analysis_data.sectionSummaries.{section}": response_data["summary"]}}
#         )

#         return jsonify(response_data)
#     except Exception as e:
#         print(f"Section summary error for {section}: {str(e)}")
#         return jsonify({"error": f"Section summary failed: {str(e)}"}), 500


# # Helper function to calculate image entropy
# def calculate_entropy(image):
#     # Convert image to grayscale if not already
#     if len(image.shape) == 3:
#         image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
#     # Calculate histogram
#     hist = cv2.calcHist([image], [0], None, [256], [0, 256])
#     hist = hist / hist.sum()  # Normalize
#     # Calculate entropy
#     entropy = stats.entropy(hist, base=2)
#     return entropy[0]

# @app.route("/api/extract-images", methods=["POST"])
# @jwt_required()
# def extract_images():
#     if "file" not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files["file"]
#     if not file.filename.endswith(".pdf"):
#         return jsonify({"error": "Only PDF files are supported"}), 400

#     try:
#         figures = []
#         pdf_bytes = file.read()
#         # Open PDF with pdfplumber to analyze text content
#         with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
#             pages = convert_from_bytes(pdf_bytes, dpi=200)  # Reduced DPI to save memory
#             logger.info(f"Converted PDF to {len(pages)} images")

#             for page_num, page_image in enumerate(pages, 1):
#                 # Check if the page contains "References" section
#                 page = pdf.pages[page_num - 1]
#                 page_text = page.extract_text() or ""
#                 if "References" in page_text or "REFERENCES" in page_text:
#                     logger.info(f"Skipping page {page_num}: Contains References section")
#                     continue

#                 img_array = np.array(page_image.convert("RGB"))
#                 img = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
#                 gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

#                 # Apply edge detection to better identify image regions
#                 edges = cv2.Canny(gray, 50, 150)
#                 # Dilate edges to connect nearby edges
#                 kernel = np.ones((5, 5), np.uint8)
#                 dilated = cv2.dilate(edges, kernel, iterations=2)
#                 contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

#                 img_idx = 0
#                 for contour in contours:
#                     x, y, w, h = cv2.boundingRect(contour)
#                     area = w * h
#                     aspect_ratio = w / h if h > 0 else 0

#                     # Enhanced filtering:
#                     # - Larger minimum area to exclude small text regions
#                     # - Exclude regions that are too large (likely entire page)
#                     # - Exclude narrow regions typical of text (aspect ratio too high or too low)
#                     # - Exclude regions with low entropy (likely text)
#                     if (area < 100000 or  # Increased minimum area
#                         area > (img.shape[0] * img.shape[1] * 0.6) or  # Reduced max area threshold
#                         aspect_ratio < 0.2 or aspect_ratio > 5):  # Aspect ratio filter
#                         continue

#                     # Extract the region and calculate entropy
#                     figure_img = img[y:y+h, x:x+w]
#                     entropy = calculate_entropy(figure_img)
#                     if entropy < 3:  # Low entropy indicates text-like content
#                         logger.info(f"Page {page_num}, Region at ({x}, {y}): Skipped, low entropy ({entropy})")
#                         continue

#                     img_idx += 1
#                     figure_img_rgb = cv2.cvtColor(figure_img, cv2.COLOR_BGR2RGB)
#                     pil_img = Image.fromarray(figure_img_rgb)
#                     buffer = io.BytesIO()
#                     pil_img.save(buffer, format="PNG")
#                     img_data = buffer.getvalue()
#                     base64_data = base64.b64encode(img_data).decode("utf-8")
#                     img_type = "png"

#                     caption = f"Figure {img_idx} on Page {page_num}"
#                     caption_match = re.search(r'(?i)Figure\s+\d+\s*[:\.\-\s]*(.*?)(?=\n|$)', page_text, re.MULTILINE)
#                     if caption_match:
#                         caption = caption_match.group(0).strip()
#                     else:
#                         lines = page_text.split('\n')
#                         for line in lines:
#                             if re.match(r'(?i)Figure\s+\d+', line):
#                                 caption = line.strip()
#                                 break

#                     figures.append({
#                         "page": page_num,
#                         "caption": caption,
#                         "image": base64_data,
#                         "type": img_type
#                     })
#                     logger.info(f"Page {page_num}, Image {img_idx}: Extracted {img_type} image, Base64 length: {len(base64_data)}, Caption: {caption}, Entropy: {entropy}")

#         logger.info(f"Total valid figures extracted: {len(figures)}")
#         return jsonify({"figures": figures})
#     except Exception as e:
#         logger.error(f"Image extraction error: {str(e)}")
#         return jsonify({"error": f"Image extraction failed: {str(e)}"}), 500


# # PERFECT
# @app.route("/api/extract-tables", methods=["POST"])
# @jwt_required()
# def extract_tables():
#     if "file" not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files["file"]
#     if not file.filename.endswith(".pdf"):
#         return jsonify({"error": "Only PDF files are supported"}), 400

#     try:
#         tables = []
#         with pdfplumber.open(file) as pdf:
#             for page_num, page in enumerate(pdf.pages, 1):
#                 page_text = page.extract_text() or ""
#                 page_tables = page.extract_tables()
#                 print(f"Page {page_num}: Found {len(page_tables)} tables")
#                 for table_idx, table in enumerate(page_tables):
#                     try:
#                         cleaned_table = [[cell or "" for cell in row] for row in table if any(cell for cell in row)]
#                         if not cleaned_table or len(cleaned_table) < 2:
#                             print(f"Page {page_num}, Table {table_idx}: Skipped, empty or invalid")
#                             continue

#                         caption = "No caption available"
#                         caption_match = re.search(r'(?i)Table\s+\d+\s*[:\.\-\s]*(.*?)(?=\n|$)', page_text, re.MULTILINE)
#                         if caption_match:
#                             caption = caption_match.group(0).strip()
#                         else:
#                             lines = page_text.split('\n')
#                             for line in lines:
#                                 if re.match(r'(?i)Table\s+\d+', line):
#                                     caption = line.strip()
#                                     break

#                         tables.append({
#                             "page": page_num,
#                             "caption": caption,
#                             "rows": cleaned_table
#                         })
#                         print(f"Page {page_num}, Table {table_idx}: Extracted, {len(cleaned_table)} rows, Caption: {caption}")
#                     except Exception as e:
#                         print(f"Page {page_num}, Table {table_idx}: Extraction failed: {str(e)}")
#         print(f"Total tables extracted: {len(tables)}")
#         return jsonify({"tables": tables})
#     except Exception as e:
#         print(f"Table extraction error: {str(e)}")
#         return jsonify({"error": f"Table extraction failed: {str(e)}"}), 500


# @app.route("/api/recommend", methods=["POST"])
# @jwt_required()
# def recommend():
#     data = request.get_json()
#     summary = data.get("summary", "")
#     if not summary:
#         return jsonify({"error": "No summary provided"}), 400

#     try:
#         prompt = f"""
#         Based on the following summary of a research paper, recommend exactly 5 recent research papers (published within the last 3 years, i.e., 2022 or later) from open-access repositories (e.g., arXiv, PubMed Central, or other public repositories). 
#         The recommendations should be closely related to the topics, methods, or findings described in the summary, such as natural language processing, PDF analysis, recommendation systems, or related fields.
#         Return the response as a JSON array of 5 objects, each with the following structure:
#         {{
#             "title": "string",
#             "published": "string (YYYY-MM-DD format)",
#             "arxiv_id": "string (or other identifier if not from arXiv)",
#             "link": "string (URL to the paper)"
#         }}
#         - Ensure the papers are from open-access sources.
#         - Prioritize recent papers (2022 or later) to ensure relevance and timeliness.
#         - Provide accurate and relevant recommendations based on your knowledge.
#         - If exact publication dates or IDs are unavailable, use reasonable estimates (e.g., "2023-01-01" for recent papers) and note any assumptions.
#         - Ensure the response is valid JSON and contains exactly 5 recommendations.
#         Summary:
#         {summary[:1000]}
#         """
#         response_text = query_gemini(prompt, max_tokens=1000)
#         logger.info(f"Raw response for recommendations: {response_text}")

#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

#         try:
#             recommendations = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             logger.error(f"JSON parsing error for recommendations: {str(e)}")
#             return jsonify({"error": f"Failed to parse response as JSON: {str(e)}"}), 500

#         if not isinstance(recommendations, list) or len(recommendations) != 5:
#             logger.error(f"Invalid number of recommendations: {len(recommendations)}")
#             return jsonify({"error": "Failed to generate exactly 5 recommendations"}), 500

#         for rec in recommendations:
#             if not all(key in rec for key in ["title", "published", "arxiv_id", "link"]):
#                 logger.error(f"Invalid recommendation structure: {rec}")
#                 return jsonify({"error": "Invalid recommendation structure"}), 500

#         logger.info(f"Successfully generated 5 recommendations")
#         return jsonify(recommendations)
#     except Exception as e:
#         logger.error(f"Recommendations error: {str(e)}")
#         return jsonify({"error": f"Recommendations failed: {str(e)}"}), 500


# @app.route("/api/chat", methods=["POST"])
# @jwt_required()
# def chat():
#     data = request.get_json()
#     user_message = data.get("message", "")
#     session_id = data.get("session_id", "")
#     if not user_message:
#         return jsonify({"error": "No message provided"}), 400
#     if not session_id:
#         return jsonify({"error": "Missing session ID"}), 400

#     try:
#         email = get_jwt_identity()
#         analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
#         if not analysis or not analysis.get("analysis_data", {}).get("text"):
#             return jsonify({"error": "No PDF context available for this session ID."}), 400

#         context = analysis["analysis_data"]["text"][:45000]
#         chat_history = analysis.get("chat_history", [])
#         logger.info(f"Chat history for session {session_id}: {chat_history}")
#         history_prompt = "\n".join([f"User: {msg['user']}\nAssistant: {msg['assistant']}" for msg in chat_history[-3:]])

#     #     prompt = f"""
#     #    You are an assistant that answers questions strictly based on the content of a research paper provided as context.
#     #     Do not use external knowledge, make assumptions, or include information beyond the provided text. 
#     #     Search the entire context, including all sections (e.g., Abstract, Introduction, Literature Survey, System Architecture, etc.), to find relevant information. 
#     #     If the question cannot be answered based on the context, respond with "This information is not available in the paper." 
#     #     If context is not available but a word from the keywords is asked to be explained, provide a concise 2-3 line explanation using your own knowledge, limited strictly to the keyword itself, and not beyond that.
#     #     Provide concise and accurate answers, directly addressing the user's question.
        
#         prompt = f"""
#         You are an assistant designed to answer questions strictly based on the content of a research paper provided as context. Your responses must adhere to the following rules:
#         1. **Context Locking**: Answer questions only using the provided paper context. Do not use external knowledge, make assumptions, or include information beyond the context.
#         2. **Comprehensive Search**: Search all sections of the paper (e.g., Abstract, Introduction, Related Work, Literature Review, Methodology, Results, Discussion, Conclusion, Future Work) to find relevant information for the user's question. Look for keywords or phrases related to the question (e.g., for "future scope," search for "future work," "next steps," "future scope").
#         3. **Specific Questions**:
#    - For questions about conclusions, look for a "Conclusion" section or summarize key findings from "Results" or "Discussion" sections.
#    - For literature gaps, search for terms like "gap," "limitation," "challenge," or sections like "Related Work" or "Literature Review."
#    - For future scope, search for terms like "future work," "next steps," or "future scope" in sections like "Discussion" or "Conclusion."
#    - For counting images, count references to "Fig.," "Figure," "Diagram," or similar terms indicating images or illustrations.
#         4. **Keyword Explanations**: If the user asks about a keyword mentioned in the paper's context or keywords list, provide a concise 2-3 sentence explanation based on the context. If the keyword is not mentioned, respond with "This information is not available in the paper."
#         5. **Fallback Response**: If the question cannot be answered based on the context, respond with "This information is not available in the paper."
#         6. **Conciseness**: Provide clear, accurate, and concise answers, directly addressing the user's question.
#         7. **Child-Friendly Explanation**: If the user asks to "explain the paper like a 10-year-old," simplify the paper's main idea and findings using age-appropriate language. Use relatable analogies (e.g., compare concepts to toys, games, or school activities), avoid technical jargon unless simplified (e.g., explain "quantum computer" as "a super cool computer that solves puzzles really fast"), and make the explanation engaging.

#         Context (up to 45,000 characters):
#         {context}
        
#         Recent conversation:
#         {history_prompt}
        
#         User question: {user_message}
        
#         Answer:
#         """
#         print(f"Chat prompt length: {len(prompt)} chars, Context length: {len(context)} chars")
#         response_text = query_gemini(prompt, max_tokens=500)

#         chat_history.append({"user": user_message, "assistant": response_text})
#         analyses_collection.update_one(
#             {"session_id": session_id, "email": email},
#             {"$set": {"chat_history": chat_history}},
#             upsert=True
#         )
#         logger.info(f"Updated chat history for session {session_id}: {chat_history}")

#         return jsonify({"response": response_text, "session_id": session_id})
#     except Exception as e:
#         logger.error(f"Chat error: {str(e)}")
#         return jsonify({"error": f"Chat failed: {str(e)}"}), 500

# @app.route("/api/paper-analysis", methods=["POST"])
# @jwt_required()
# def save_analysis():
#     data = request.get_json()
#     email = get_jwt_identity()
#     analysis_data = data.get("analysis_data")
#     if not analysis_data:
#         return jsonify({"error": "Missing analysis data"}), 400

#     try:
#         recommendations = analysis_data.get("recommendations", [])
#         logger.info(f"Initial recommendations count: {len(recommendations)}")
        
#         unique_recommendations = []
#         seen_titles = set()
#         for rec in recommendations:
#             title = rec.get("title", "")
#             if title and title not in seen_titles:
#                 seen_titles.add(title)
#                 unique_recommendations.append(rec)
#             if len(unique_recommendations) >= 5:
#                 break

#         logger.info(f"After deduplication, recommendations count: {len(unique_recommendations)}")

#         analysis_doc = {
#             "email": email,
#             "session_id": analysis_data.get("session_id", str(time.time())),
#             "analysis_data": analysis_data,
#             "tables": analysis_data.get("tables", []),
#             "figures": analysis_data.get("figures", []),
#             "recommendations": unique_recommendations,
#             "chat_history": analysis_data.get("messages", []),
#             "created_at": time.time()
#         }
#         logger.info(f"Saving analysis for {email}, recommendations: {len(unique_recommendations)}, chat_history: {analysis_doc['chat_history']}")
#         result = analyses_collection.insert_one(analysis_doc)
#         logger.info(f"Stored analysis for {email}: {result.inserted_id}")
#         return jsonify({"message": "Analysis saved", "id": str(result.inserted_id)})
#     except Exception as e:
#         logger.error(f"Save analysis error: {str(e)}")
#         return jsonify({"error": f"Failed to save analysis: {str(e)}"}), 500

# @app.route("/api/paper-analysis/<analysis_id>", methods=["DELETE"])
# @jwt_required()
# def delete_analysis(analysis_id):
#     email = get_jwt_identity()
#     try:
#         result = analyses_collection.delete_one({"_id": ObjectId(analysis_id), "email": email})
#         if result.deleted_count == 1:
#             logger.info(f"Deleted analysis {analysis_id} for {email}")
#             return jsonify({"message": "Analysis deleted successfully"}), 200
#         else:
#             return jsonify({"error": "Analysis not found or unauthorized"}), 404
#     except Exception as e:
#         logger.error(f"Delete analysis error: {str(e)}")
#         return jsonify({"error": f"Failed to delete analysis: {str(e)}"}), 500

# @app.route("/api/user-data", methods=["GET"])
# @jwt_required()
# def user_data():
#     email = get_jwt_identity()
#     try:
#         analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
#         for analysis in analyses:
#             analysis["_id"] = str(analysis["_id"])
#         logger.info(f"Retrieved {len(analyses)} analyses for {email}")
#         return jsonify({"recent_analysis": analyses})
#     except Exception as e:
#         logger.error(f"User data error: {str(e)}")
#         return jsonify({"error": f"Failed to retrieve user data: {str(e)}"}), 500

# @app.route("/api/generate-pdf", methods=["POST"])
# @jwt_required()
# def generate_pdf():
#     try:
#         data = request.get_json()
#         analysis_data = data.get("analysis_data")
#         if not analysis_data:
#             return jsonify({"error": "Missing analysis data"}), 400

#         logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
#         logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

#         prompt = f"""
#         Generate a professional PDF report for a research paper based on the following analysis data.
#         Format the content as a structured JSON object with sections suitable for a conference presentation.
#         Include the following sections:
#         - Title: The paper's title
#         - Authors: List of authors
#         - Abstract: A concise summary (use the provided summary)
#         - Methodology: A concise summary of the paper's methodology (use sectionSummaries.Methodology if available, otherwise infer from the text)
#         - Keywords: List of keywords
#         - Tables: Full table data including captions and rows/columns (limit to 3 tables for brevity)
#         - Figures: Full figure data including captions (limit to 3 figures for brevity, exclude image data)
#         - Citations: Key citations (limit to 5 for brevity, preserve original numbering like "[1]")
#         - Recommendations: All recommended papers (up to 5, include title, published date, and link)
#         For tables, include each table's caption and its full data as a list of lists (rows and columns).
#         For figures, include each figure's caption and page number (e.g., "Figure 1 on Page 2").
#         For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
#         For citations, preserve the original numbering (e.g., "[1]") as it appears in the analysis data.
#         For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
#         Return the response as a valid JSON object:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "abstract": "string",
#             "methodology": "string",
#             "keywords": ["string", ...],
#             "tables": [
#                 {{
#                     "caption": "string",
#                     "rows": [["string", ...], ...]
#                 }},
#                 ...
#             ],
#             "figures": [
#                 {{
#                     "caption": "string",
#                     "page": number
#                 }},
#                 ...
#             ],
#             "citations": ["string", ...],
#             "recommendations": ["string", ...]
#         }}
#         Analysis Data:
#         {json.dumps(analysis_data, indent=2)[:10000]}
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
#         report_content = json.loads(response_text)

#         logger.info(f"report_content.recommendations: {report_content.get('recommendations', [])}")
#         logger.info(f"report_content.methodology: {report_content.get('methodology', 'Not provided')}")

#         buffer = BytesIO()
#         doc = SimpleDocTemplate(buffer, pagesize=letter)
#         styles = getSampleStyleSheet()
#         title_style = ParagraphStyle(
#             name='TitleStyle',
#             fontSize=16,
#             leading=20,
#             alignment=1,
#             spaceAfter=12,
#             fontName='Helvetica-Bold'
#         )
#         heading_style = ParagraphStyle(
#             name='HeadingStyle',
#             fontSize=12,
#             leading=14,
#             spaceAfter=8,
#             fontName='Helvetica-Bold'
#         )
#         normal_style = ParagraphStyle(
#             name='NormalStyle',
#             fontSize=10,
#             leading=12,
#             spaceAfter=6
#         )

#         story = []
#         story.append(Paragraph(report_content.get("title", "Untitled"), title_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Authors", heading_style))
#         authors = ", ".join(report_content.get("authors", ["No authors found"]))
#         story.append(Paragraph(authors, normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Abstract", heading_style))
#         story.append(Paragraph(report_content.get("abstract", "No summary available"), normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Methodology", heading_style))
#         methodology = report_content.get("methodology", "No methodology available")
#         story.append(Paragraph(methodology, normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Keywords", heading_style))
#         keywords = ", ".join(report_content.get("keywords", []))
#         story.append(Paragraph(keywords, normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Tables", heading_style))
#         tables = report_content.get("tables", [])[:3]
#         if tables:
#             for i, table in enumerate(tables, 1):
#                 caption = table.get("caption", f"Table {i}")
#                 rows = table.get("rows", [])
#                 if rows:
#                     table_data = [[Paragraph(str(cell), normal_style) for cell in row] for row in rows]
#                     pdf_table = Table(table_data)
#                     pdf_table.setStyle(TableStyle([
#                         ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
#                         ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
#                         ('VALIGN', (0, 0), (-1, -1), 'TOP'),
#                         ('FONTSIZE', (0, 0), (-1, -1), 8),
#                     ]))
#                     story.append(Paragraph(caption, normal_style))
#                     story.append(pdf_table)
#                     story.append(Spacer(1, 12))
#                 else:
#                     story.append(Paragraph(f"{caption}: No data available", normal_style))
#         else:
#             story.append(Paragraph("No tables available", normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Figures", heading_style))
#         figures = report_content.get("figures", [])[:3]
#         if figures:
#             for figure in figures:
#                 caption = f"{figure.get('caption', 'Figure')} on Page {figure.get('page', 'Unknown')}"
#                 story.append(Paragraph(caption, normal_style))
#         else:
#             story.append(Paragraph("No figures available", normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Citations", heading_style))
#         citations = report_content.get("citations", [])[:5]
#         for citation in citations:
#             story.append(Paragraph(citation, normal_style))
#         if not citations:
#             story.append(Paragraph("No citations available", normal_style))
#         story.append(Spacer(1, 12))

#         story.append(Paragraph("Recommendations", heading_style))
#         recommendations = report_content.get("recommendations", [])
#         if recommendations:
#             for i, rec in enumerate(recommendations, 1):
#                 story.append(Paragraph(f"{i}. {rec}", normal_style))
#         else:
#             story.append(Paragraph("No recommendations available", normal_style))
#         story.append(Spacer(1, 12))

#         doc.build(story)
#         buffer.seek(0)
#         return send_file(
#             buffer,
#             as_attachment=True,
#             download_name=f"{report_content.get('title', 'report')}.pdf",
#             mimetype='application/pdf'
#         )
#     except Exception as e:
#         logger.error(f"PDF generation error: {str(e)}")
#         return jsonify({"error": f"Failed to generate PDF: {str(e)}"}), 500

# @app.route("/api/generate-ppt", methods=["POST"])
# @jwt_required()
# def generate_ppt():
#     try:
#         data = request.get_json()
#         analysis_data = data.get("analysis_data")
#         if not analysis_data:
#             return jsonify({"error": "Missing analysis data"}), 400

#         logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
#         logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

#         prompt = f"""
#         Generate content for a PowerPoint presentation for a research paper based on the following analysis data.
#         Format the content as a structured JSON object with sections suitable for a conference presentation.
#         Include the following sections for slides:
#         - Title: The paper's title and authors
#         - Summary: A concise summary (use the provided summary)
#         - Keywords: List of keywords
#         - Tables: Full table data including captions and rows/columns (limit to 2 tables for brevity)
#         - Figures: Full figure data including captions and page numbers (limit to 2 figures for brevity, exclude image data)
#         - Citations: Key citations (limit to 3, preserve original numbering like "[1]")
#         - Recommendations: All recommended papers (up to 5, include title, published date, and link)
#         For tables, include each table's caption and its full data as a list of lists (rows and columns).
#         For figures, include each figure's caption and page number (e.g., "Figure 1 on Page 2").
#         For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
#         For citations, preserve the original numbering (e.g., "[1]") as it appears in the analysis data.
#         For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
#         Return the response as a valid JSON object:
#         {{
#             "title": "string",
#             "authors": ["string", ...],
#             "summary": "string",
#             "methodology": "string",
#             "keywords": ["string", ...],
#             "tables": [
#                 {{
#                     "caption": "string",
#                     "rows": [["string", ...], ...]
#                 }},
#                 ...
#             ],
#             "figures": [
#                 {{
#                     "caption": "string",
#                     "page": number
#                 }},
#                 ...
#             ],
#             "citations": ["string", ...],
#             "recommendations": ["string", ...]
#         }}
#         Analysis Data:
#         {json.dumps(analysis_data, indent=2)[:10000]}
#         """
#         response_text = query_gemini(prompt, max_tokens=4000)
#         response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
#         ppt_content = json.loads(response_text)

#         logger.info(f"ppt_content.recommendations: {ppt_content.get('recommendations', [])}")

#         prs = Presentation()
#         slide_width = prs.slide_width
#         slide_height = prs.slide_height

#         slide_layout = prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         subtitle = slide.placeholders[1]
#         title.text = ppt_content.get("title", "Untitled")
#         title.text_frame.paragraphs[0].font.size = Pt(32)
#         subtitle.text = ", ".join(ppt_content.get("authors", ["No authors found"]))
#         subtitle.text_frame.paragraphs[0].font.size = Pt(20)

#         slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Summary"
#         content.text = ppt_content.get("summary", "No summary available")
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         # slide = prs.slides.add_slide(slide_layout)
#         # title = slide.shapes.title
#         # content = slide.placeholders[1]
#         # title.text = "Methodology"
#         # methodology = ppt_content.get("methodology", "No methodology available")
#         # if not methodology or methodology.strip() == "":
#         #     methodology = "No methodology available"
#         #     logger.warning("Methodology field empty or missing in ppt_content")
#         # content.text = methodology
#         # content.text_frame.paragraphs[0].font.size = Pt(16)
#         # for p in content.text_frame.paragraphs:
#         #     p.font.size = Pt(16)

#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Keywords"
#         keywords = ppt_content.get("keywords", [])
#         content.text = ", ".join(keywords) if keywords else "No keywords available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         tables = ppt_content.get("tables", [])[:2]
#         if tables:
#             for i, table in enumerate(tables, 1):
#                 slide = prs.slides.add_slide(slide_layout)
#                 title = slide.shapes.title
#                 title.text = f"Table {i}: {table.get('caption', 'Table')}"

#                 rows = table.get("rows", [])
#                 if rows:
#                     num_rows = min(len(rows), 10)
#                     num_cols = min(len(rows[0]), 6) if rows else 0
#                     if num_rows > 0 and num_cols > 0:
#                         table_shape = slide.shapes.add_table(
#                             rows=num_rows,
#                             cols=num_cols,
#                             left=Inches(1),
#                             top=Inches(1.5),
#                             width=Inches(8),
#                             height=Inches(4)
#                         )
#                         ppt_table = table_shape.table
#                         for row_idx in range(num_rows):
#                             for col_idx in range(num_cols):
#                                 cell = ppt_table.cell(row_idx, col_idx)
#                                 cell.text = str(rows[row_idx][col_idx])
#                                 cell.text_frame.paragraphs[0].font.size = Pt(12)
#                         ppt_table.first_row = True
#                     else:
#                         content = slide.placeholders[1]
#                         content.text = "Table data too large or invalid"
#                         content.text_frame.paragraphs[0].font.size = Pt(18)
#                 else:
#                     content = slide.placeholders[1]
#                     content.text = "No table data available"
#                     content.text_frame.paragraphs[0].font.size = Pt(18)
#         else:
#             slide = prs.slides.add_slide(slide_layout)
#             title = slide.shapes.title
#             content = slide.placeholders[1]
#             title.text = "Tables"
#             content.text = "No tables available"
#             content.text_frame.paragraphs[0].font.size = Pt(18)

#         figures = ppt_content.get("figures", [])[:2]
#         if figures:
#             slide = prs.slides.add_slide(slide_layout)
#             title = slide.shapes.title
#             content = slide.placeholders[1]
#             title.text = "Figures"
#             for figure in figures:
#                 p = content.text_frame.add_paragraph()
#                 p.text = f"{figure.get('caption', 'Figure')} on Page {figure.get('page', 'Unknown')}"
#                 p.font.size = Pt(16)
#         else:
#             slide = prs.slides.add_slide(slide_layout)
#             title = slide.shapes.title
#             content = slide.placeholders[1]
#             title.text = "Figures"
#             content.text = "No figures available"
#             content.text_frame.paragraphs[0].font.size = Pt(18)

#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Citations"
#         citations = ppt_content.get("citations", [])[:3]
#         if citations:
#             for citation in citations:
#                 p = content.text_frame.add_paragraph()
#                 p.text = citation
#                 p.font.size = Pt(16)
#         else:
#             content.text = "No citations available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Recommendations"
#         recommendations = ppt_content.get("recommendations", [])
#         if recommendations:
#             for i, rec in enumerate(recommendations, 1):
#                 p = content.text_frame.add_paragraph()
#                 p.text = f"{i}. {rec}"
#                 p.font.size = Pt(14)
#                 p.level = 1
#         else:
#             content.text = "No recommendations available"
#         content.text_frame.paragraphs[0].font.size = Pt(18)

#         slide = prs.slides.add_slide(slide_layout)
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = "Thank You"
#         content.text = "Questions?"
#         content.text_frame.paragraphs[0].font.size = Pt(24)
#         content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

#         buffer = BytesIO()
#         prs.save(buffer)
#         buffer.seek(0)
#         return send_file(
#             buffer,
#             as_attachment=True,
#             download_name=f"{ppt_content.get('title', 'presentation')}.pptx",
#             mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
#         )
#     except Exception as e:
#         logger.error(f"PPT generation error: {str(e)}")
#         return jsonify({"error": f"Failed to generate PPT: {str(e)}"}), 500

# @app.route("/api/convert-citations", methods=["POST"])
# @jwt_required()
# def convert_citations():
#     try:
#         data = request.get_json()
#         citations = data.get("citations", [])
#         style = data.get("style", "APA").upper()

#         if not citations:
#             return jsonify({"error": "No citations provided"}), 400
#         if style not in ["APA", "MLA", "CHICAGO", "IEEE"]:
#             return jsonify({"error": f"Invalid citation style: {style}. Supported styles: APA, MLA, Chicago, IEEE"}), 400

#         prompt = f"""
#         Convert the following citations to the {style} citation style.
#         The input citations are raw strings extracted from a research paper's References section.
#         Parse each citation to identify key components (e.g., authors, title, journal, year, volume, pages, DOI, URL) and reformat them according to the {style} style guidelines.
#         - For APA: Follow the 7th edition (e.g., Author, A. A. (Year). Title. Journal, Volume(Issue), pages. DOI).
#         - For MLA: Follow the 9th edition (e.g., Author. "Title." Journal, Volume, Issue, Year, pages, DOI or URL).
#         - For Chicago (Author-Date): Follow the 17th edition (e.g., Author, First. Year. "Title." Journal Volume (Issue): pages. DOI).
#         - For IEEE: Follow the IEEE style (e.g., [1] A. Author, "Title," Journal, vol. Volume, no. Issue, pp. pages, Year).
#         - If a citation lacks sufficient information (e.g., missing year or title), include as much as possible and note "[Incomplete]" at the end.
#         - Preserve all available metadata (e.g., DOIs, URLs) and include them as required by the style.
#         - Return the response as a JSON array of strings, where each string is a citation formatted in the {style} style.
#         - Ensure the response is valid JSON and matches the number of input citations.
#         - Wrap the JSON in triple backticks with 'json' identifier:
#         ```json
#         ["formatted citation 1", "formatted citation 2", ...]
#         ```
#         Input Citations:
#         {json.dumps(citations, indent=2)[:10000]}
#         """
#         max_tokens = 3000
#         response_text = query_gemini(prompt, max_tokens=max_tokens)
#         logger.info(f"Raw response for {style} conversion: {response_text}")

#         json_match = re.search(r'```json\n([\s\S]*?)\n```', response_text, re.MULTILINE)
#         if json_match:
#             response_text = json_match.group(1).strip()
#         else:
#             response_text = re.sub(r'^.*?\[|\].*?$', '', response_text, flags=re.DOTALL).strip()
#             response_text = f'[{response_text}]'

#         try:
#             converted_citations = json.loads(response_text)
#         except json.JSONDecodeError as e:
#             logger.error(f"JSON parsing error for citation conversion: {str(e)}")
#             logger.error(f"Problematic response_text: {response_text}")
#             converted_citations = [f"{citation} [Conversion failed]" for citation in citations]
#             logger.warning(f"Falling back to original citations with error note")

#         if not isinstance(converted_citations, list) or len(converted_citations) != len(citations):
#             logger.error(f"Invalid converted citations: Expected {len(citations)}, got {len(converted_citations)}")
#             converted_citations = [f"{citation} [Invalid conversion]" for citation in citations]

#         logger.info(f"Successfully converted {len(converted_citations)} citations to {style} style")
#         return jsonify({"converted_citations": converted_citations})
#     except Exception as e:
#         logger.error(f"Citation conversion error: {str(e)}")
#         return jsonify({"error": f"Failed to convert citations: {str(e)}"}), 500


# @app.route("/api/feedback", methods=["POST"])
# @jwt_required()
# def submit_feedback():
#     try:
#         email = get_jwt_identity()
#         data = request.get_json()
#         session_id = data.get("session_id", "")
#         feedback = data.get("feedback", "").lower()

#         if not session_id:
#             logger.error("Missing session_id in feedback request")
#             return jsonify({"error": "Missing session ID"}), 400
#         if feedback not in ["up", "down"]:
#             logger.error(f"Invalid feedback value: {feedback}")
#             return jsonify({"error": "Feedback must be 'up' or 'down'"}), 400

#         analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
#         if not analysis:
#             logger.error(f"No analysis found for session_id: {session_id}, email: {email}")
#             return jsonify({"error": "No analysis found for this session ID"}), 404

#         analyses_collection.update_one(
#             {"session_id": session_id, "email": email},
#             {"$set": {"feedback": feedback}},
#             upsert=True
#         )
#         logger.info(f"Feedback submitted for session {session_id} by {email}: {feedback}")
#         return jsonify({"message": "Feedback submitted successfully"}), 200
#     except Exception as e:
#         logger.error(f"Feedback submission error: {str(e)}")
#         return jsonify({"error": f"Failed to submit feedback: {str(e)}"}), 500

# if __name__ == "__main__":
#     app.run(debug=True, host="0.0.0.0", port=5000)

















from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import pdfplumber
import base64
import os
import json
import google.generativeai as genai
from dotenv import load_dotenv
import re
import time
from google.api_core.exceptions import ResourceExhausted
import requests
from sentence_transformers import SentenceTransformer
from numpy import dot
from numpy.linalg import norm
import imghdr
import pymongo
import scipy.stats as stats
from bson.objectid import ObjectId
import logging
from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity, set_access_cookies, unset_jwt_cookies
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pdf2image import convert_from_bytes
import numpy as np
import cv2
from PIL import Image
import io

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "http://localhost:5173", "methods": ["GET", "POST", "OPTIONS", "DELETE"], "allow_headers": ["Content-Type", "Authorization"]}})

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
MONGO_URI = os.getenv("MONGO_URI")
JWT_SECRET_KEY = os.getenv("JWT_SECRET_KEY", "your-secret-key")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found in .env file")
if not MONGO_URI:
    raise ValueError("MONGO_URI not found in .env file")
genai.configure(api_key=GOOGLE_API_KEY)

# Initialize MongoDB
try:
    client = pymongo.MongoClient(MONGO_URI)
    db = client.get_database("researchXtract")
    analyses_collection = db.analyses
    users_collection = db.users
    logger.info("Connected to MongoDB Atlas")
except Exception as e:
    logger.error(f"Failed to connect to MongoDB: {str(e)}")
    raise

# JWT Configuration
app.config["JWT_SECRET_KEY"] = JWT_SECRET_KEY
app.config["JWT_TOKEN_LOCATION"] = ["headers", "cookies"]
jwt = JWTManager(app)

# Handle OPTIONS requests for CORS preflight
@app.route('/api/<path:path>', methods=['OPTIONS'])
def handle_options(path):
    return '', 200

# Initialize Gemini model
MODEL_NAME = "gemini-2.0-flash-lite"
model = genai.GenerativeModel(MODEL_NAME)

# Initialize sentence transformer for semantic similarity
sentence_model = SentenceTransformer('all-MiniLM-L6-v2')

# Store PDF text and chat history in memory (per session)
pdf_context = {"text": "", "session_id": None}
chat_history = []

def validate_input(text):
    return len(text) > 0 and len(text) < 250000

def validate_research_paper(pdf_file):
    """
    Validates if the uploaded PDF is a research paper by checking for common research paper features.
    Returns True if valid, False otherwise.
    """
    try:
        pdf_file.seek(0)  # Reset file pointer
        with pdfplumber.open(pdf_file) as pdf:
            text = "".join(page.extract_text() or "" for page in pdf.pages)
        
        if not text:
            logger.info("PDF is empty or contains no extractable text")
            return False

        # Convert text to lowercase for case-insensitive matching
        text_lower = text.lower()

        # Define criteria for research papers
        score = 0
        required_score = 3  # Minimum score to qualify as a research paper

        # Check for common sections
        sections = ["abstract", "introduction", "references", "bibliography", "conclusion", "methodology", "results", "discussion"]
        for section in sections:
            if section in text_lower:
                score += 1
                logger.info(f"Found section: {section}")

        # Check for citations (e.g., [1], (Author, Year))
        citation_patterns = [
            r'\[\d+\]',  # e.g., [1], [12]
            r'\(\w+,\s*\d{4}\)',  # e.g., (Smith, 2020)
            r'\(\w+ et al\.,\s*\d{4}\)'  # e.g., (Smith et al., 2020)
        ]
        for pattern in citation_patterns:
            if re.search(pattern, text):
                score += 1
                logger.info(f"Found citation pattern: {pattern}")
                break

        # Check for keywords section or technical terms
        if "keywords" in text_lower or "key words" in text_lower:
            score += 1
            logger.info("Found keywords section")

        # Check for author list (e.g., names followed by affiliations or commas)
        author_pattern = r'(\w+\s+\w+\s*,?\s*){2,}'  # At least two names
        if re.search(author_pattern, text, re.MULTILINE):
            score += 1
            logger.info("Found potential author list")

        # Check for title-like text on the first page (bold or large font not easily detectable, so use heuristic)
        first_page = pdf.pages[0].extract_text() or ""
        if len(first_page.split('\n')[0].strip()) > 10:  # Assume title is longer than 10 chars
            score += 1
            logger.info("Found potential title on first page")

        logger.info(f"Research paper validation score: {score}/{required_score}")
        return score >= required_score

    except Exception as e:
        logger.error(f"Validation error: {str(e)}")
        return False

def query_gemini(prompt, max_tokens=512, retries=3, delay=10, image_data=None):
    for attempt in range(retries):
        try:
            if image_data:
                logger.warning("Image data provided but not supported in this implementation.")
            response = model.generate_content(prompt, generation_config={"max_output_tokens": max_tokens})
            return response.text.strip()
        except ResourceExhausted:
            if attempt < retries - 1:
                print(f"Quota exceeded, retrying in {delay} seconds... (Attempt {attempt + 1}/{retries})")
                time.sleep(delay)
                continue
            raise Exception("Network issue")
        except Exception:
            raise Exception("An error occurred")

def cosine_similarity(a, b):
    return dot(a, b) / (norm(a) * norm(b))

# Authentication Routes
@app.route("/api/register", methods=["POST"])
def register():
    logger.info("Received POST request to /api/register")
    data = request.get_json()
    email = data.get("email")
    password = data.get("password")
    name = data.get("name", "Unknown")
    if not email or not password:
        return jsonify({"error": "Email and password are required"}), 400
    if users_collection.find_one({"email": email}):
        return jsonify({"error": "Email already registered"}), 400
    user_data = {"email": email, "password": password, "name": name, "token": None}
    logger.info(f"Registering user: {user_data}")
    users_collection.insert_one(user_data)
    return jsonify({"message": "User registered successfully."}), 201

@app.route("/api/login", methods=["POST"])
def login():
    logger.info("Received POST request to /api/login")
    data = request.get_json()
    email = data.get("email")
    password = data.get("password")
    user = users_collection.find_one({"email": email, "password": password})
    if not user:
        logger.error(f"Login failed for {email}: Invalid credentials")
        return jsonify({"error": "Invalid credentials"}), 401
    access_token = create_access_token(identity=email)
    users_collection.update_one({"email": email}, {"$set": {"token": access_token}})
    user_data = {"email": user["email"], "name": user.get("name", "Unknown")}
    logger.info(f"Login successful for {email}, user data: {user_data}")
    response = jsonify({"access_token": access_token, "user": user_data})
    set_access_cookies(response, access_token)
    return response, 200

@app.route("/api/user", methods=["GET"])
@jwt_required()
def get_user():
    email = get_jwt_identity()
    user = users_collection.find_one({"email": email}, {"_id": 0, "password": 0})
    if not user:
        logger.error(f"User not found for email: {email}")
        return jsonify({"error": "User not found"}), 404
    logger.info(f"User data for {email}: {user}")
    return jsonify({"user": user}), 200

@app.route("/api/logout", methods=["POST"])
@jwt_required()
def logout():
    logger.info("Received POST request to /api/logout")
    email = get_jwt_identity()
    users_collection.update_one({"email": email}, {"$set": {"token": None}})
    response = jsonify({"message": "Logged out successfully"})
    unset_jwt_cookies(response)
    return response, 200

@app.route("/api/user-history", methods=["GET"])
@jwt_required()
def user_history():
    email = get_jwt_identity()
    try:
        analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
        for analysis in analyses:
            analysis["_id"] = str(analysis["_id"])
            if "recommendations" in analysis and len(analysis["recommendations"]) > 5:
                analysis["recommendations"] = analysis["recommendations"][:5]
                logger.info(f"Limited recommendations for analysis {analysis['_id']} to 5")
            logger.info(f"History item {analysis['_id']}: chat_history = {analysis.get('chat_history', [])}")
        return jsonify({"history": analyses}), 200
    except Exception:
        logger.error("An error occurred while retrieving user history")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/analyze-pdf", methods=["POST"])
@jwt_required()
def analyze_pdf():

    time.sleep(20)

    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files["file"]
    word_limit = request.form.get("word_limit", 150)
    if not file.filename.endswith(".pdf"):
        return jsonify({"error": "Only PDF files are supported"}), 400

    # Validate if the PDF is a research paper
    if not validate_research_paper(file):
        logger.error("Uploaded file is not a valid research paper")
        return jsonify({"error": "Invalid document: Only research papers are accepted"}), 400

    try:
        file.seek(0)  # Reset file pointer after validation
        with pdfplumber.open(file) as pdf:
            text = "".join(page.extract_text() or "" for page in pdf.pages)
        print(f"Extracted text length: {len(text)}")
        print(f"First 500 chars of extracted text: {text[:500]}")
        if not validate_input(text):
            return jsonify({"error": "Invalid input: text length must be 1-249999 characters"}), 400

        global pdf_context, chat_history
        session_id = str(time.time())
        pdf_context = {"text": text, "session_id": session_id}
        chat_history = []

        prompt_text = text[:45000]
        if len(text) > 45000:
            ref_start = text.rfind("REFERENCES") or text.rfind("References")
            if ref_start != -1:
                ref_text = text[ref_start:ref_start+20000]
                prompt_text = text[:25000] + "\n" + ref_text
                prompt_text = prompt_text[:45000]
            print(f"Text truncated to {len(prompt_text)} chars, References included: {ref_start != -1}")

        prompt = f"""
        Analyze the following research paper text and extract the specified fields.
        Return the response as a JSON object with the following structure:
        {{
            "title": "string",
            "authors": ["string", ...],
            "summary": "string (summarize in {word_limit} words)",
            "keywords": ["string", ...],
            "citations": ["string", ...]
        }}
        - For "title", extract the exact title of the paper, typically found at the top of the first page.
        - For "authors", list all author names as they appear (e.g., "Christian Szegedy").
        - For "summary", provide a concise summary of the paper's abstract and introduction in {word_limit} words, formatted as concise bullet points.
        - For "keywords", extract all keywords listed in the paper's "Keywords" section. If no "Keywords" section exists, infer 5–10 keywords from the abstract and introduction, focusing on technical terms and concepts central to the paper (e.g., "Inception", "convolutional neural networks").
        - For "citations", extract ALL citations listed in the paper's "References" or bibliography section, preserving their exact formatting as they appear in the original text (e.g., "[1] S. Arora, A. Bhaskara, ..."). Include the numbering (e.g., "[1]", "[2]") as it appears in the paper, along with all other details such as author names, titles, journal names, years, DOIs, URLs, and any LaTeX formatting or special characters. Do not modify, reformat, or remove any part of the citation, including the numbering.
        - If a field cannot be extracted, use appropriate defaults (e.g., "Untitled" for title, [] for lists, "No summary available" for summary).
        - Ensure the response is valid JSON.
        Text (up to 45,000 characters):
        {prompt_text}
        """
        response_text = query_gemini(prompt, max_tokens=4000)
        print(f"Raw response: {response_text}")

        response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            print(f"JSON parsing error")
            return jsonify({"error": "An error occurred"}), 500

        analysis_data = {
            "title": result.get("title", "Untitled"),
            "authors": result.get("authors", ["No authors found"]),
            "summary": result.get("summary", "No summary available"),
            "keywords": result.get("keywords", []),
            "citations": result.get("citations", []),
            "session_id": session_id,
            "text": text
        }
        print(f"Extracted keywords: {analysis_data['keywords']}")
        print(f"Extracted citations (with numbering check): {analysis_data['citations'][:5]}")
        return jsonify(analysis_data)
    except Exception:
        print(f"Analyze PDF error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/summarize-section", methods=["POST"])
@jwt_required()
def summarize_section():
    data = request.get_json()
    section = data.get("section", "")
    session_id = data.get("session_id", "")
    word_limit = 150

    if not section:
        return jsonify({"error": "No section provided"}), 400
    if not session_id:
        return jsonify({"error": "Missing session ID"}), 400

    try:
        email = get_jwt_identity()
        analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
        if not analysis or not analysis.get("analysis_data", {}).get("text"):
            return jsonify({"error": "No PDF context available for this session ID."}), 400

        prompt_text = analysis["analysis_data"]["text"][:45000]
        print(f"Summarizing section: {section}, Context length: {len(prompt_text)} chars")

        prompt = f"""
        Analyze the following research paper text and provide a concise summary of the specified section in {word_limit} words.
        - Section to summarize: "{section}"
        - If the section is "Entire Paper", summarize the abstract and introduction (as in the default summary).
        - For other sections (e.g., "Introduction", "Literature Review", "Methodology", "Results", "Conclusion"), identify the section by its title or content (e.g., "Related Work" for "Literature Review", "Experiments" for "Results").
        - If the section is not found, return: "This section is not available in the paper."
        - Return the response as a JSON object:
        {{
            "section": "{section}",
            "summary": "string (summarize in {word_limit} words)"
        }}
        - Ensure the summary is concise, accurate, and based only on the provided text formatted as concise bullet points.
        - Ensure the response is valid JSON.
        Text (up to 45,000 characters):
        {prompt_text}
        """
        response_text = query_gemini(prompt, max_tokens=1000)
        print(f"Raw response for section {section}: {response_text}")

        response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            print(f"JSON parsing error for section {section}")
            return jsonify({"error": "An error occurred"}), 500

        response_data = {
            "section": result.get("section", section),
            "summary": result.get("summary", "This section is not available in the paper.")
        }
        print(f"Section summary for {section}: {response_data['summary'][:100]}...")

        analyses_collection.update_one(
            {"session_id": session_id, "email": email},
            {"$set": {f"analysis_data.sectionSummaries.{section}": response_data["summary"]}}
        )

        return jsonify(response_data)
    except Exception:
        print(f"Section summary error for {section}")
        return jsonify({"error": "An error occurred"}), 500

# Helper function to calculate image entropy
def calculate_entropy(image):
    if len(image.shape) == 3:
        image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    hist = cv2.calcHist([image], [0], None, [256], [0, 256])
    hist = hist / hist.sum()
    entropy = stats.entropy(hist, base=2)
    return entropy[0]

@app.route("/api/extract-images", methods=["POST"])
@jwt_required()
def extract_images():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files["file"]
    if not file.filename.endswith(".pdf"):
        return jsonify({"error": "Only PDF files are supported"}), 400

    # Validate if the PDF is a research paper
    if not validate_research_paper(file):
        logger.error("Uploaded file is not a valid research paper")
        return jsonify({"error": "Invalid document: Only research papers are accepted"}), 400

    try:
        figures = []
        file.seek(0)  # Reset file pointer after validation
        pdf_bytes = file.read()
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            pages = convert_from_bytes(pdf_bytes, dpi=200)
            logger.info(f"Converted PDF to {len(pages)} images")

            for page_num, page_image in enumerate(pages, 1):
                page = pdf.pages[page_num - 1]
                page_text = page.extract_text() or ""
                if "References" in page_text or "REFERENCES" in page_text:
                    logger.info(f"Skipping page {page_num}: Contains References section")
                    continue

                img_array = np.array(page_image.convert("RGB"))
                img = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

                edges = cv2.Canny(gray, 50, 150)
                kernel = np.ones((5, 5), np.uint8)
                dilated = cv2.dilate(edges, kernel, iterations=2)
                contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

                img_idx = 0
                for contour in contours:
                    x, y, w, h = cv2.boundingRect(contour)
                    area = w * h
                    aspect_ratio = w / h if h > 0 else 0

                    if (area < 100000 or
                        area > (img.shape[0] * img.shape[1] * 0.6) or
                        aspect_ratio < 0.2 or aspect_ratio > 5):
                        continue

                    figure_img = img[y:y+h, x:x+w]
                    entropy = calculate_entropy(figure_img)
                    if entropy < 3:
                        logger.info(f"Page {page_num}, Region at ({x}, {y}): Skipped, low entropy ({entropy})")
                        continue

                    img_idx += 1
                    figure_img_rgb = cv2.cvtColor(figure_img, cv2.COLOR_BGR2RGB)
                    pil_img = Image.fromarray(figure_img_rgb)
                    buffer = io.BytesIO()
                    pil_img.save(buffer, format="PNG")
                    img_data = buffer.getvalue()
                    base64_data = base64.b64encode(img_data).decode("utf-8")
                    img_type = "png"

                    caption = f"Figure {img_idx} on Page {page_num}"
                    caption_match = re.search(r'(?i)Figure\s+\d+\s*[:\.\-\s]*(.*?)(?=\n|$)', page_text, re.MULTILINE)
                    if caption_match:
                        caption = caption_match.group(0).strip()
                    else:
                        lines = page_text.split('\n')
                        for line in lines:
                            if re.match(r'(?i)Figure\s+\d+', line):
                                caption = line.strip()
                                break

                    figures.append({
                        "page": page_num,
                        "caption": caption,
                        "image": base64_data,
                        "type": img_type
                    })
                    logger.info(f"Page {page_num}, Image {img_idx}: Extracted {img_type} image, Base64 length: {len(base64_data)}, Caption: {caption}, Entropy: {entropy}")

        logger.info(f"Total valid figures extracted: {len(figures)}")
        return jsonify({"figures": figures})
    except Exception:
        logger.error("Image extraction error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/extract-tables", methods=["POST"])
@jwt_required()
def extract_tables():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files["file"]
    if not file.filename.endswith(".pdf"):
        return jsonify({"error": "Only PDF files are supported"}), 400

    # Validate if the PDF is a research paper
    if not validate_research_paper(file):
        logger.error("Uploaded file is not a valid research paper")
        return jsonify({"error": "Invalid document: Only research papers are accepted"}), 400

    try:
        tables = []
        file.seek(0)  # Reset file pointer after validation
        with pdfplumber.open(file) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                page_tables = page.extract_tables()
                print(f"Page {page_num}: Found {len(page_tables)} tables")
                for table_idx, table in enumerate(page_tables):
                    try:
                        cleaned_table = [[cell or "" for cell in row] for row in table if any(cell for cell in row)]
                        if not cleaned_table or len(cleaned_table) < 2:
                            print(f"Page {page_num}, Table {table_idx}: Skipped, empty or invalid")
                            continue

                        caption = "No caption available"
                        caption_match = re.search(r'(?i)Table\s+\d+\s*[:\.\-\s]*(.*?)(?=\n|$)', page_text, re.MULTILINE)
                        if caption_match:
                            caption = caption_match.group(0).strip()
                        else:
                            lines = page_text.split('\n')
                            for line in lines:
                                if re.match(r'(?i)Table\s+\d+', line):
                                    caption = line.strip()
                                    break

                        tables.append({
                            "page": page_num,
                            "caption": caption,
                            "rows": cleaned_table
                        })
                        print(f"Page {page_num}, Table {table_idx}: Extracted, {len(cleaned_table)} rows, Caption: {caption}")
                    except Exception:
                        print(f"Page {page_num}, Table {table_idx}: Extraction failed")
        print(f"Total tables extracted: {len(tables)}")
        return jsonify({"tables": tables})
    except Exception:
        print("Table extraction error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/recommend", methods=["POST"])
@jwt_required()
def recommend():
    data = request.get_json()
    summary = data.get("summary", "")
    if not summary:
        return jsonify({"error": "No summary provided"}), 400

    try:
        prompt = f"""
        Based on the following summary of a research paper, recommend exactly 5 recent research papers (published within the last 3 years, i.e., 2022 or later) from open-access repositories (e.g., arXiv, PubMed Central, or other public repositories). 
        The recommendations should be closely related to the topics, methods, or findings described in the summary, such as natural language processing, PDF analysis, recommendation systems, or related fields.
        Return the response as a JSON array of 5 objects, each with the following structure:
        {{
            "title": "string",
            "published": "string (YYYY-MM-DD format)",
            "arxiv_id": "string (or other identifier if not from arXiv)",
            "link": "string (URL to the paper)"
        }}
        - Ensure the papers are from open-access sources.
        - Prioritize recent papers (2022 or later) to ensure relevance and timeliness.
        - Provide accurate and relevant recommendations based on your knowledge.
        - If exact publication dates or IDs are unavailable, use reasonable estimates (e.g., "2023-01-01" for recent papers) and note any assumptions.
        - Ensure the response is valid JSON and contains exactly 5 recommendations.
        Summary:
        {summary[:1000]}
        """
        response_text = query_gemini(prompt, max_tokens=1000)
        logger.info(f"Raw response for recommendations: {response_text}")

        response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()

        try:
            recommendations = json.loads(response_text)
        except json.JSONDecodeError:
            logger.error("JSON parsing error for recommendations")
            return jsonify({"error": "An error occurred"}), 500

        if not isinstance(recommendations, list) or len(recommendations) != 5:
            logger.error(f"Invalid number of recommendations: {len(recommendations)}")
            return jsonify({"error": "An error occurred"}), 500

        for rec in recommendations:
            if not all(key in rec for key in ["title", "published", "arxiv_id", "link"]):
                logger.error(f"Invalid recommendation structure: {rec}")
                return jsonify({"error": "An error occurred"}), 500

        logger.info(f"Successfully generated 5 recommendations")
        return jsonify(recommendations)
    except Exception:
        logger.error("Recommendations error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/chat", methods=["POST"])
@jwt_required()
def chat():
    data = request.get_json()
    user_message = data.get("message", "")
    session_id = data.get("session_id", "")
    if not user_message:
        return jsonify({"error": "No message provided"}), 400
    if not session_id:
        return jsonify({"error": "Missing session ID"}), 400

    try:
        email = get_jwt_identity()
        analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
        if not analysis or not analysis.get("analysis_data", {}).get("text"):
            return jsonify({"error": "No PDF context available for this session ID."}), 400

        context = analysis["analysis_data"]["text"][:45000]
        chat_history = analysis.get("chat_history", [])
        logger.info(f"Chat history for session {session_id}: {chat_history}")
        history_prompt = "\n".join([f"User: {msg['user']}\nAssistant: {msg['assistant']}" for msg in chat_history[-3:]])

        prompt = f"""
        You are an assistant designed to answer questions strictly based on the content of a research paper provided as context. Your responses must adhere to the following rules:
        1. **Context Locking**: Answer questions only using the provided paper context. Do not use external knowledge, make assumptions, or include information beyond the context.
        2. **Comprehensive Search**: Search all sections of the paper (e.g., Abstract, Introduction, Related Work, Literature Review, Methodology, Results, Discussion, Conclusion, Future Work) to find relevant information for the user's question. Look for keywords or phrases related to the question (e.g., for "future scope," search for "future work," "next steps," "future scope").
        3. **Specific Questions**:
           - For questions about conclusions, look for a "Conclusion" section or summarize key findings from "Results" or "Discussion" sections.
           - For literature gaps, search for terms like "gap," "limitation," "challenge," or sections like "Related Work" or "Literature Review."
           - For future scope, search for terms like "future work," "next steps," or "future scope" in sections like "Discussion" or "Conclusion."
           - For counting images, count references to "Fig.," "Figure," "Diagram," or similar terms indicating images or illustrations.
        4. **Keyword Explanations**: If the user asks about a keyword mentioned in the paper's context or keywords list, provide a concise 2-3 sentence explanation based on the context. If the keyword is not mentioned, respond with "This information is not available in the paper."
        5. **Fallback Response**: If the question cannot be answered based on the context, respond with "This information is not available in the paper."
        6. **Conciseness**: Provide clear, accurate, and concise answers, directly addressing the user's question.
        7. **Child-Friendly Explanation**: If the user asks to "explain the paper like a 10-year-old," simplify the paper's main idea and findings using age-appropriate language. Use relatable analogies (e.g., compare concepts to toys, games, or school activities), avoid technical jargon unless simplified (e.g., explain "quantum computer" as "a super cool computer that solves puzzles really fast"), and make the explanation engaging.

        Context (up to 45,000 characters):
        {context}
        
        Recent conversation:
        {history_prompt}
        
        User question: {user_message}
        
        Answer:
        """
        print(f"Chat prompt length: {len(prompt)} chars, Context length: {len(context)} chars")
        response_text = query_gemini(prompt, max_tokens=500)

        chat_history.append({"user": user_message, "assistant": response_text})
        analyses_collection.update_one(
            {"session_id": session_id, "email": email},
            {"$set": {"chat_history": chat_history}},
            upsert=True
        )
        logger.info(f"Updated chat history for session {session_id}: {chat_history}")

        return jsonify({"response": response_text, "session_id": session_id})
    except Exception:
        logger.error("Chat error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/paper-analysis", methods=["POST"])
@jwt_required()
def save_analysis():
    data = request.get_json()
    email = get_jwt_identity()
    analysis_data = data.get("analysis_data")
    if not analysis_data:
        return jsonify({"error": "Missing analysis data"}), 400

    try:
        recommendations = analysis_data.get("recommendations", [])
        logger.info(f"Initial recommendations count: {len(recommendations)}")
        
        unique_recommendations = []
        seen_titles = set()
        for rec in recommendations:
            title = rec.get("title", "")
            if title and title not in seen_titles:
                seen_titles.add(title)
                unique_recommendations.append(rec)
            if len(unique_recommendations) >= 5:
                break

        logger.info(f"After deduplication, recommendations count: {len(unique_recommendations)}")

        analysis_doc = {
            "email": email,
            "session_id": analysis_data.get("session_id", str(time.time())),
            "analysis_data": analysis_data,
            "tables": analysis_data.get("tables", []),
            "figures": analysis_data.get("figures", []),
            "recommendations": unique_recommendations,
            "chat_history": analysis_data.get("messages", []),
            "created_at": time.time()
        }
        logger.info(f"Saving analysis for {email}, recommendations: {len(unique_recommendations)}, chat_history: {analysis_doc['chat_history']}")
        result = analyses_collection.insert_one(analysis_doc)
        logger.info(f"Stored analysis for {email}: {result.inserted_id}")
        return jsonify({"message": "Analysis saved", "id": str(result.inserted_id)})
    except Exception:
        logger.error("Save analysis error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/paper-analysis/<analysis_id>", methods=["DELETE"])
@jwt_required()
def delete_analysis(analysis_id):
    email = get_jwt_identity()
    try:
        result = analyses_collection.delete_one({"_id": ObjectId(analysis_id), "email": email})
        if result.deleted_count == 1:
            logger.info(f"Deleted analysis {analysis_id} for {email}")
            return jsonify({"message": "Analysis deleted successfully"}), 200
        else:
            return jsonify({"error": "Analysis not found or unauthorized"}), 404
    except Exception:
        logger.error("Delete analysis error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/user-data", methods=["GET"])
@jwt_required()
def user_data():
    email = get_jwt_identity()
    try:
        analyses = list(analyses_collection.find({"email": email}).sort("created_at", -1).limit(10))
        for analysis in analyses:
            analysis["_id"] = str(analysis["_id"])
        logger.info(f"Retrieved {len(analyses)} analyses for {email}")
        return jsonify({"recent_analysis": analyses})
    except Exception:
        logger.error("User data error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/generate-pdf", methods=["POST"])
@jwt_required()
def generate_pdf():
    try:
        data = request.get_json()
        analysis_data = data.get("analysis_data")
        if not analysis_data:
            return jsonify({"error": "Missing analysis data"}), 400

        logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
        logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

        prompt = f"""
        Generate a professional PDF report for a research paper based on the following analysis data.
        Format the content as a structured JSON object with sections suitable for a conference presentation.
        Include the following sections:
        - Title: The paper's title
        - Authors: List of authors
        - Abstract: A concise summary (use the provided summary)
        - Methodology: A concise summary of the paper's methodology (use sectionSummaries.Methodology if available, otherwise infer from the text)
        - Keywords: List of keywords
        - Tables: Full table data including captions and rows/columns (limit to 3 tables for brevity)
        - Figures: Full figure data including captions (limit to 3 figures for brevity, exclude image data)
        - Citations: Key citations (limit to 5 for brevity, preserve original numbering like "[1]")
        - Recommendations: All recommended papers (up to 5, include title, published date, and link)
        For tables, include each table's caption and its full data as a list of lists (rows and columns).
        For figures, include each figure's caption and page number (e.g., "Figure 1 on Page 2").
        For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
        For citations, preserve the original numbering (e.g., "[1]") as it appears in the analysis data.
        For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
        Return the response as a valid JSON object:
        {{
            "title": "string",
            "authors": ["string", ...],
            "abstract": "string",
            "methodology": "string",
            "keywords": ["string", ...],
            "tables": [
                {{
                    "caption": "string",
                    "rows": [["string", ...], ...]
                }},
                ...
            ],
            "figures": [
                {{
                    "caption": "string",
                    "page": number
                }},
                ...
            ],
            "citations": ["string", ...],
            "recommendations": ["string", ...]
        }}
        Analysis Data:
        {json.dumps(analysis_data, indent=2)[:10000]}
        """
        response_text = query_gemini(prompt, max_tokens=4000)
        response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
        report_content = json.loads(response_text)

        logger.info(f"report_content.recommendations: {report_content.get('recommendations', [])}")
        logger.info(f"report_content.methodology: {report_content.get('methodology', 'Not provided')}")

        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            name='TitleStyle',
            fontSize=16,
            leading=20,
            alignment=1,
            spaceAfter=12,
            fontName='Helvetica-Bold'
        )
        heading_style = ParagraphStyle(
            name='HeadingStyle',
            fontSize=12,
            leading=14,
            spaceAfter=8,
            fontName='Helvetica-Bold'
        )
        normal_style = ParagraphStyle(
            name='NormalStyle',
            fontSize=10,
            leading=12,
            spaceAfter=6
        )

        story = []
        story.append(Paragraph(report_content.get("title", "Untitled"), title_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Authors", heading_style))
        authors = ", ".join(report_content.get("authors", ["No authors found"]))
        story.append(Paragraph(authors, normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Abstract", heading_style))
        story.append(Paragraph(report_content.get("abstract", "No summary available"), normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Methodology", heading_style))
        methodology = report_content.get("methodology", "No methodology available")
        story.append(Paragraph(methodology, normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Keywords", heading_style))
        keywords = ", ".join(report_content.get("keywords", []))
        story.append(Paragraph(keywords, normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Tables", heading_style))
        tables = report_content.get("tables", [])[:3]
        if tables:
            for i, table in enumerate(tables, 1):
                caption = table.get("caption", f"Table {i}")
                rows = table.get("rows", [])
                if rows:
                    table_data = [[Paragraph(str(cell), normal_style) for cell in row] for row in rows]
                    pdf_table = Table(table_data)
                    pdf_table.setStyle(TableStyle([
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ]))
                    story.append(Paragraph(caption, normal_style))
                    story.append(pdf_table)
                    story.append(Spacer(1, 12))
                else:
                    story.append(Paragraph(f"{caption}: No data available", normal_style))
        else:
            story.append(Paragraph("No tables available", normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Figures", heading_style))
        figures = report_content.get("figures", [])[:3]
        if figures:
            for figure in figures:
                caption = f"{figure.get('caption', 'Figure')} on Page {figure.get('page', 'Unknown')}"
                story.append(Paragraph(caption, normal_style))
        else:
            story.append(Paragraph("No figures available", normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Citations", heading_style))
        citations = report_content.get("citations", [])[:5]
        for citation in citations:
            story.append(Paragraph(citation, normal_style))
        if not citations:
            story.append(Paragraph("No citations available", normal_style))
        story.append(Spacer(1, 12))

        story.append(Paragraph("Recommendations", heading_style))
        recommendations = report_content.get("recommendations", [])
        if recommendations:
            for i, rec in enumerate(recommendations, 1):
                story.append(Paragraph(f"{i}. {rec}", normal_style))
        else:
            story.append(Paragraph("No recommendations available", normal_style))
        story.append(Spacer(1, 12))

        doc.build(story)
        buffer.seek(0)
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"{report_content.get('title', 'report')}.pdf",
            mimetype='application/pdf'
        )
    except Exception:
        logger.error("PDF generation error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/generate-ppt", methods=["POST"])
@jwt_required()
def generate_ppt():
    try:
        data = request.get_json()
        analysis_data = data.get("analysis_data")
        if not analysis_data:
            return jsonify({"error": "Missing analysis data"}), 400

        logger.info(f"analysis_data.recommendations: {analysis_data.get('recommendations', [])}")
        logger.info(f"analysis_data.sectionSummaries.Methodology: {analysis_data.get('sectionSummaries', {}).get('Methodology', 'Not provided')}")

        prompt = f"""
        Generate content for a PowerPoint presentation for a research paper based on the following analysis data.
        Format the content as a structured JSON object with sections suitable for a conference presentation.
        Include the following sections for slides:
        - Title: The paper's title and authors
        - Summary: A concise summary (use the provided summary)
        - Keywords: List of keywords
        - Tables: Full table data including captions and rows/columns (limit to 2 tables for brevity)
        - Figures: Full figure data including captions and page numbers (limit to 2 figures for brevity, exclude image data)
        - Citations: Key citations (limit to 3, preserve original numbering like "[1]")
        - Recommendations: All recommended papers (up to 5, include title, published date, and link)
        For tables, include each table's caption and its full data as a list of lists (rows and columns).
        For figures, include each figure's caption and page number (e.g., "Figure 1 on Page 2").
        For methodology, provide a 100-150 word summary of the paper's methodology section or equivalent (e.g., "Methods", "Approach"). If unavailable, return "No methodology available."
        For citations, preserve the original numbering (e.g., "[1]") as it appears in the analysis data.
        For recommendations, format each as a string combining title, published date, and link (e.g., "Title (Published: YYYY-MM-DD) - Link: URL"). Use all recommendations provided in the analysis data.
        Return the response as a valid JSON object:
        {{
            "title": "string",
            "authors": ["string", ...],
            "summary": "string",
            "methodology": "string",
            "keywords": ["string", ...],
            "tables": [
                {{
                    "caption": "string",
                    "rows": [["string", ...], ...]
                }},
                ...
            ],
            "figures": [
                {{
                    "caption": "string",
                    "page": number
                }},
                ...
            ],
            "citations": ["string", ...],
            "recommendations": ["string", ...]
        }}
        Analysis Data:
        {json.dumps(analysis_data, indent=2)[:10000]}
        """
        response_text = query_gemini(prompt, max_tokens=4000)
        response_text = re.sub(r'^```json\n|\n```$', '', response_text).strip()
        ppt_content = json.loads(response_text)

        logger.info(f"ppt_content.recommendations: {ppt_content.get('recommendations', [])}")

        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = ppt_content.get("title", "Untitled")
        title.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text = ", ".join(ppt_content.get("authors", ["No authors found"]))
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Summary"
        content.text = ppt_content.get("summary", "No summary available")
        content.text_frame.paragraphs[0].font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Keywords"
        keywords = ppt_content.get("keywords", [])
        content.text = ", ".join(keywords) if keywords else "No keywords available"
        content.text_frame.paragraphs[0].font.size = Pt(18)

        tables = ppt_content.get("tables", [])[:2]
        if tables:
            for i, table in enumerate(tables, 1):
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Table {i}: {table.get('caption', 'Table')}"

                rows = table.get("rows", [])
                if rows:
                    num_rows = min(len(rows), 10)
                    num_cols = min(len(rows[0]), 6) if rows else 0
                    if num_rows > 0 and num_cols > 0:
                        table_shape = slide.shapes.add_table(
                            rows=num_rows,
                            cols=num_cols,
                            left=Inches(1),
                            top=Inches(1.5),
                            width=Inches(8),
                            height=Inches(4)
                        )
                        ppt_table = table_shape.table
                        for row_idx in range(num_rows):
                            for col_idx in range(num_cols):
                                cell = ppt_table.cell(row_idx, col_idx)
                                cell.text = str(rows[row_idx][col_idx])
                                cell.text_frame.paragraphs[0].font.size = Pt(12)
                        ppt_table.first_row = True
                    else:
                        content = slide.placeholders[1]
                        content.text = "Table data too large or invalid"
                        content.text_frame.paragraphs[0].font.size = Pt(18)
                else:
                    content = slide.placeholders[1]
                    content.text = "No table data available"
                    content.text_frame.paragraphs[0].font.size = Pt(18)
        else:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Tables"
            content.text = "No tables available"
            content.text_frame.paragraphs[0].font.size = Pt(18)

        figures = ppt_content.get("figures", [])[:2]
        if figures:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Figures"
            for figure in figures:
                p = content.text_frame.add_paragraph()
                p.text = f"{figure.get('caption', 'Figure')} on Page {figure.get('page', 'Unknown')}"
                p.font.size = Pt(16)
        else:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Figures"
            content.text = "No figures available"
            content.text_frame.paragraphs[0].font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Citations"
        citations = ppt_content.get("citations", [])[:3]
        if citations:
            for citation in citations:
                p = content.text_frame.add_paragraph()
                p.text = citation
                p.font.size = Pt(16)
        else:
            content.text = "No citations available"
        content.text_frame.paragraphs[0].font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Recommendations"
        recommendations = ppt_content.get("recommendations", [])
        if recommendations:
            for i, rec in enumerate(recommendations, 1):
                p = content.text_frame.add_paragraph()
                p.text = f"{i}. {rec}"
                p.font.size = Pt(14)
                p.level = 1
        else:
            content.text = "No recommendations available"
        content.text_frame.paragraphs[0].font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Thank You"
        content.text = "Questions?"
        content.text_frame.paragraphs[0].font.size = Pt(24)
        content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"{ppt_content.get('title', 'presentation')}.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception:
        logger.error("PPT generation error")
        return jsonify({"error": "An error occurred"}), 500


@app.route("/api/convert-citations", methods=["POST"])
@jwt_required()
def convert_citations():
    try:
        data = request.get_json()
        citations = data.get("citations", [])
        style = data.get("style", "APA").upper()

        if not citations:
            return jsonify({"error": "No citations provided"}), 400
        if style not in ["APA", "MLA", "CHICAGO", "IEEE"]:
            return jsonify({"error": f"Invalid citation style: {style}. Supported styles: APA, MLA, Chicago, IEEE"}), 400

        prompt = f"""
        Convert the following citations to the {style} citation style.
        The input citations are raw strings extracted from a research paper's References section.
        Parse each citation to identify key components (e.g., authors, title, journal, year, volume, pages, DOI, URL) and reformat them according to the {style} style guidelines.
        - For APA: Follow the 7th edition (e.g., Author, A. A. (Year). Title. Journal, Volume(Issue), pages.).
        - For MLA: Follow the 9th edition (e.g., Author. "Title." Journal, Volume, Issue, Year, pages, DOI or URL).
        - For Chicago (Author-Date): Follow the 17th edition (e.g., Author, First. Year. "Title." Journal Volume (Issue): pages.).
        - For IEEE: Follow the IEEE style (e.g., [1] A. Author, "Title," Journal, vol. Volume, no. Issue, pp. pages, Year).
        - If a citation lacks sufficient information (e.g., missing year or title), include as much as possible and note "[Incomplete]" at the end.
        - Preserve all available metadata (e.g., DOIs, URLs) and include them as required by the style.
        - Return the response as a JSON array of strings, where each string is a citation formatted in the {style} style.
        - Ensure the response is valid JSON and matches the number of input citations.
        - Wrap the JSON in triple backticks with 'json' identifier:
        ```json
        ["formatted citation 1", "formatted citation 2", ...]
        ```
        Input Citations:
        {json.dumps(citations, indent=2)[:10000]}
        """
        max_tokens = 1500
        response_text = query_gemini(prompt, max_tokens=max_tokens)
        logger.info(f"Raw response for {style} conversion: {response_text}")

        # Corrected regex to properly capture JSON block
        json_match = re.search(r'```json\n([\s\S]*?)\n```', response_text, re.MULTILINE)
        if json_match:
            response_text = json_match.group(1).strip()
        else:
            # Fallback: Try to extract JSON array directly
            json_start = response_text.find('[')
            json_end = response_text.rfind(']')
            if json_start != -1 and json_end != -1:
                response_text = response_text[json_start:json_end + 1]
            else:
                response_text = f'[{response_text}]'

        try:
            converted_citations = json.loads(response_text)
        except json.JSONDecodeError:
            logger.error("JSON parsing error for citation conversion")
            converted_citations = [f"{citation} [Conversion failed]" for citation in citations]
            logger.warning(f"Falling back to original citations with error note")

        if not isinstance(converted_citations, list) or len(converted_citations) != len(citations):
            logger.error(f"Invalid converted citations: Expected {len(citations)}, got {len(converted_citations)}")
            converted_citations = [f"{citation} [Invalid conversion]" for citation in citations]

        logger.info(f"Successfully converted {len(converted_citations)} citations to {style} style")
        return jsonify({"converted_citations": converted_citations})
    except Exception:
        logger.error("Citation conversion error")
        return jsonify({"error": "An error occurred"}), 500

@app.route("/api/feedback", methods=["POST"])
@jwt_required()
def submit_feedback():
    try:
        email = get_jwt_identity()
        data = request.get_json()
        session_id = data.get("session_id", "")
        feedback = data.get("feedback", "").lower()

        if not session_id:
            logger.error("Missing session_id in feedback request")
            return jsonify({"error": "Missing session ID"}), 400
        if feedback not in ["up", "down"]:
            logger.error(f"Invalid feedback value: {feedback}")
            return jsonify({"error": "Feedback must be 'up' or 'down'"}), 400

        analysis = analyses_collection.find_one({"session_id": session_id, "email": email})
        if not analysis:
            logger.error(f"No analysis found for session_id: {session_id}, email: {email}")
            return jsonify({"error": "No analysis found for this session ID"}), 404

        analyses_collection.update_one(
            {"session_id": session_id, "email": email},
            {"$set": {"feedback": feedback}},
            upsert=True
        )
        logger.info(f"Feedback submitted for session {session_id} by {email}: {feedback}")
        return jsonify({"message": "Feedback submitted successfully"}), 200
    except Exception:
        logger.error("Feedback submission error")
        return jsonify({"error": "An error occurred"}), 500

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)